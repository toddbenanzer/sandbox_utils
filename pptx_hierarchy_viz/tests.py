from hierarchical_visualizer import HierarchicalVisualizer
from node import Node
from pathlib import Path
from powerpoint_exporter import PowerPointExporter
from pptx import Presentation
from unittest.mock import patch, Mock
from your_module_name import fetch_data_from_api
from your_module_name import format_visual_elements
from your_module_name import load_data_from_csv, convert_to_hierarchy
from your_module_name import load_data_from_json
from your_module_name import load_data_from_xml
from your_module_name import validate_hierarchy_data
import csv
import json
import os
import pytest
import requests
import xml.etree.ElementTree as ET


@pytest.fixture
def visualizer():
    return HierarchicalVisualizer()

def test_create_tree_structure(visualizer):
    data = {'root': {'child1': {}, 'child2': {}}}
    result = visualizer.create_tree_structure(data)
    assert result == "Tree visualization object", "Tree structure creation failed."

def test_create_org_chart(visualizer):
    data = {'CEO': {'VP1': {'Manager1': {}}, 'VP2': {}}}
    result = visualizer.create_org_chart(data)
    assert result == "Org chart visualization object", "Org chart creation failed."

def test_export_to_powerpoint(visualizer, tmpdir):
    filepath = tmpdir.join("test_presentation.pptx")
    visualizer.export_to_powerpoint(str(filepath))

    # Verify PowerPoint file creation
    prs = Presentation(str(filepath))
    slide_count = len(prs.slides)
    assert slide_count == 1, "PowerPoint slide creation failed."
    assert prs.slides[0].shapes.title.text == "My Hierarchical Visualization", \
        "Title text is incorrect."

    # Clean up
    os.remove(filepath)



def test_node_initialization():
    node = Node(id=1, label="Root Node")
    assert node.id == 1
    assert node.label == "Root Node"
    assert node.parent is None
    assert node.children == []
    assert node.style == {'color': 'black', 'shape': 'rectangle', 'line_thickness': 1}

def test_node_style_initialization():
    node = Node(id=2, label="Styled Node", color="red", shape="circle", line_thickness=2)
    assert node.style['color'] == "red"
    assert node.style['shape'] == "circle"
    assert node.style['line_thickness'] == 2

def test_set_style():
    node = Node(id=3, label="Node")
    node.set_style(color="blue", shape="triangle", line_thickness=3)
    assert node.style['color'] == "blue"
    assert node.style['shape'] == "triangle"
    assert node.style['line_thickness'] == 3

def test_add_child():
    parent = Node(id=4, label="Parent")
    child = Node(id=5, label="Child")
    parent.add_child(child)
    assert len(parent.children) == 1
    assert parent.children[0] == child
    assert child.parent == parent



@pytest.fixture
def sample_template(tmpdir):
    # Create a temporary PowerPoint file
    template_path = tmpdir.join("template.pptx")
    prs = Presentation()
    prs.save(str(template_path))
    return str(template_path)

@pytest.fixture
def exporter(sample_template):
    export_instance = PowerPointExporter(slide_template=sample_template)
    export_instance.init_powerpoint()
    return export_instance

def test_init_powerpoint(sample_template):
    exporter = PowerPointExporter(slide_template=sample_template)
    assert exporter.prs is None
    exporter.init_powerpoint()
    assert isinstance(exporter.prs, Presentation)

def test_add_slide_no_init():
    exporter = PowerPointExporter(slide_template="path/to/nonexistent/template.pptx")
    with pytest.raises(ValueError, match=r"Presentation not initialized.*"):
        exporter.add_slide(visual_data="Sample Data")

def test_add_slide(exporter):
    exporter.add_slide(visual_data="Sample Visual Data")
    assert len(exporter.prs.slides) == 1
    slide = exporter.prs.slides[0]
    assert slide.shapes.title.text == "Visual Data Representation"
    text_box_content = slide.shapes[1].text_frame.text
    assert text_box_content == "Sample Visual Data"

def test_save_presentation(exporter, tmpdir):
    export_path = tmpdir.join("presentation.pptx")
    exporter.save_presentation(str(export_path))
    assert Path(export_path).exists()



def test_load_data_from_csv_file_not_found():
    with pytest.raises(FileNotFoundError):
        load_data_from_csv("nonexistent_file.csv")

def test_load_data_from_csv_invalid_format(tmpdir):
    # Create a temporary invalid CSV file
    invalid_file = tmpdir.join("invalid.csv")
    invalid_file.write("This is not a CSV\nJust some random text")
    with pytest.raises(csv.Error):
        load_data_from_csv(str(invalid_file))

def test_load_data_from_csv_empty_file(tmpdir):
    # Create an empty CSV file
    empty_file = tmpdir.join("empty.csv")
    empty_file.write("")
    with pytest.raises(ValueError, match="CSV file is empty or cannot be parsed into a hierarchy."):
        load_data_from_csv(str(empty_file))

def test_load_data_from_csv_valid_file_structure(tmpdir):
    # Create a temporary valid CSV file
    valid_file = tmpdir.join("valid.csv")
    valid_file.write("id,parent_id\n1,\n2,1\n3,1\n4,2")
    
    def mock_convert_to_hierarchy(data):
        # A basic mock to return the data as received for testing
        return {"root": data}

    # Patch the convert_to_hierarchy function to make testing easier
    original_convert = convert_to_hierarchy
    try:
        setattr(you_module_name, 'convert_to_hierarchy', mock_convert_to_hierarchy)
        expected_result = {
            "root": [
                {"id": "1", "parent_id": ""},
                {"id": "2", "parent_id": "1"},
                {"id": "3", "parent_id": "1"},
                {"id": "4", "parent_id": "2"}
            ]
        }
        result = load_data_from_csv(str(valid_file))
        assert result == expected_result
    finally:
        setattr(you_module_name, 'convert_to_hierarchy', original_convert)



def test_load_data_from_json_file_not_found():
    with pytest.raises(FileNotFoundError):
        load_data_from_json("nonexistent_file.json")

def test_load_data_from_json_invalid_json(tmpdir):
    # Create a temporary invalid JSON file
    invalid_file = tmpdir.join("invalid.json")
    invalid_file.write("This is not a JSON file")
    with pytest.raises(json.JSONDecodeError):
        load_data_from_json(str(invalid_file))

def test_load_data_from_json_non_dict_structure(tmpdir):
    # Create a JSON file with valid JSON but not a dictionary
    invalid_structure_file = tmpdir.join("invalid_structure.json")
    invalid_structure_file.write('["This", "is", "a", "list"]')
    with pytest.raises(ValueError, match="JSON content is not in the expected hierarchical format."):
        load_data_from_json(str(invalid_structure_file))

def test_load_data_from_json_valid_file(tmpdir):
    # Create a temporary valid JSON file
    valid_file = tmpdir.join("valid.json")
    valid_file.write('{"root": {"child": "value"}}')
    result = load_data_from_json(str(valid_file))
    assert result == {"root": {"child": "value"}}



def test_load_data_from_xml_file_not_found():
    with pytest.raises(FileNotFoundError):
        load_data_from_xml("nonexistent_file.xml")

def test_load_data_from_xml_invalid_xml(tmpdir):
    # Create a temporary invalid XML file
    invalid_file = tmpdir.join("invalid.xml")
    invalid_file.write("<root><unclosed></root>")
    with pytest.raises(ET.ParseError):
        load_data_from_xml(str(invalid_file))

def test_load_data_from_xml_non_dict_structure(tmpdir):
    # XML is inherently hierarchical, so this test checks empty or unexpected structures
    empty_file = tmpdir.join("empty.xml")
    empty_file.write("")
    with pytest.raises(ET.ParseError):
        load_data_from_xml(str(empty_file))

def test_load_data_from_xml_valid_file(tmpdir):
    # Create a valid XML file
    valid_file = tmpdir.join("valid.xml")
    valid_file.write("""
        <root>
            <child1 attr1="value1">
                <grandchild attr2="value2"/>
            </child1>
            <child2 attr3="value3"/>
        </root>
    """)
    
    expected_result = {
        'root': {
            'attributes': {},
            'children': [
                {'child1': {
                    'attributes': {'attr1': 'value1'},
                    'children': [
                        {'grandchild': {
                            'attributes': {'attr2': 'value2'},
                            'children': []
                        }}
                    ]
                }},
                {'child2': {
                    'attributes': {'attr3': 'value3'},
                    'children': []
                }}
            ]
        }
    }

    result = load_data_from_xml(str(valid_file))
    assert result == expected_result




def test_fetch_data_from_api_success():
    # Mock response
    example_data = {'root': {'child': 'value'}}
    with patch('requests.get') as mock_get:
        mock_response = Mock()
        mock_response.json.return_value = example_data
        mock_response.raise_for_status = Mock()
        mock_get.return_value = mock_response

        result = fetch_data_from_api("http://example.com/hierarchy")
        assert result == example_data


def test_fetch_data_from_api_http_error():
    # Mock HTTP error
    with patch('requests.get') as mock_get:
        mock_response = Mock()
        mock_response.raise_for_status.side_effect = requests.exceptions.HTTPError("HTTP Error")
        mock_get.return_value = mock_response

        with pytest.raises(requests.exceptions.RequestException, match="HTTP Error"):
            fetch_data_from_api("http://example.com/hierarchy")


def test_fetch_data_from_api_json_decode_error():
    # Mock JSON decode error
    with patch('requests.get') as mock_get:
        mock_response = Mock()
        mock_response.raise_for_status = Mock()
        mock_response.json.side_effect = json.JSONDecodeError("JSON Error", "", 0)
        mock_get.return_value = mock_response

        with pytest.raises(json.JSONDecodeError, match="JSON Error"):
            fetch_data_from_api("http://example.com/hierarchy")


def test_fetch_data_from_api_value_error():
    # Mock response with a non-dict JSON
    with patch('requests.get') as mock_get:
        mock_response = Mock()
        mock_response.raise_for_status = Mock()
        mock_response.json.return_value = ["not", "a", "dict"]
        mock_get.return_value = mock_response

        with pytest.raises(ValueError, match="JSON content is not in the expected hierarchical format."):
            fetch_data_from_api("http://example.com/hierarchy")


def test_fetch_data_from_api_request_exception():
    # Mock network request error
    with patch('requests.get', side_effect=requests.exceptions.RequestException("Network Error")):
        with pytest.raises(requests.exceptions.RequestException, match="Network Error"):
            fetch_data_from_api("http://example.com/hierarchy")



def test_validate_hierarchy_data_valid_structure():
    valid_data = {
        'root': {
            'attributes': {},
            'children': [
                {
                    'child1': {
                        'attributes': {'attr1': 'value1'},
                        'children': []
                    }
                }
            ]
        }
    }
    assert validate_hierarchy_data(valid_data) is True

def test_validate_hierarchy_data_missing_keys():
    missing_keys_data = {
        'root': {
            'attributes': {}
            # 'children' key is missing
        }
    }
    assert validate_hierarchy_data(missing_keys_data) is False

def test_validate_hierarchy_data_children_not_list():
    children_not_list_data = {
        'root': {
            'attributes': {},
            'children': {}  # children should be a list
        }
    }
    assert validate_hierarchy_data(children_not_list_data) is False

def test_validate_hierarchy_data_non_dict_input():
    non_dict_data = "This is not a dictionary"
    assert validate_hierarchy_data(non_dict_data) is False

def test_validate_hierarchy_data_circular_reference():
    # Circular reference setup
    circular_node = {
        'attributes': {},
        'children': []
    }
    circular_node['children'].append({'circular_child': circular_node})  # Adds self-reference

    circular_data = {
        'root': circular_node
    }
    assert validate_hierarchy_data(circular_data) is False




def test_format_visual_elements_success():
    style_config = {
        'node_color': 'blue',
        'node_shape': 'circle',
        'line_thickness': 2,
        'font_size': 10,
        'font_color': 'black'
    }
    expected_output = {
        'node_color': 'blue',
        'node_shape': 'circle',
        'line_thickness': 2,
        'font_size': 10,
        'font_color': 'black'
    }
    assert format_visual_elements(style_config) == expected_output


def test_format_visual_elements_missing_keys():
    style_config = {
        'node_color': 'blue',
        'node_shape': 'circle',
        'line_thickness': 2,
        'font_size': 10
        # Missing 'font_color'
    }
    with pytest.raises(KeyError, match="Missing required style attribute: font_color"):
        format_visual_elements(style_config)


def test_format_visual_elements_invalid_node_color_type():
    style_config = {
        'node_color': 123,  # Should be a string
        'node_shape': 'circle',
        'line_thickness': 2,
        'font_size': 10,
        'font_color': 'black'
    }
    with pytest.raises(TypeError, match="Expected a string for 'node_color'"):
        format_visual_elements(style_config)


def test_format_visual_elements_invalid_shape_type():
    style_config = {
        'node_color': 'blue',
        'node_shape': 123,  # Should be a string
        'line_thickness': 2,
        'font_size': 10,
        'font_color': 'black'
    }
    with pytest.raises(TypeError, match="Expected a string for 'node_shape'"):
        format_visual_elements(style_config)


def test_format_visual_elements_invalid_line_thickness_type():
    style_config = {
        'node_color': 'blue',
        'node_shape': 'circle',
        'line_thickness': 'thick',  # Should be a number
        'font_size': 10,
        'font_color': 'black'
    }
    with pytest.raises(TypeError, match="Expected a number for 'line_thickness'"):
        format_visual_elements(style_config)


def test_format_visual_elements_invalid_font_size_type():
    style_config = {
        'node_color': 'blue',
        'node_shape': 'circle',
        'line_thickness': 2,
        'font_size': 'small',  # Should be an integer
        'font_color': 'black'
    }
    with pytest.raises(TypeError, match="Expected an integer for 'font_size'"):
        format_visual_elements(style_config)


def test_format_visual_elements_invalid_font_color_type():
    style_config = {
        'node_color': 'blue',
        'node_shape': 'circle',
        'line_thickness': 2,
        'font_size': 10,
        'font_color': 123  # Should be a string
    }
    with pytest.raises(TypeError, match="Expected a string for 'font_color'"):
        format_visual_elements(style_config)
