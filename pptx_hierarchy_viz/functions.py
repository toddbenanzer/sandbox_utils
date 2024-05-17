
from pptx import Presentation
from pptx.util import Inches, Pt, RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
import csv
import pandas as pd


def create_blank_slide(prs):
    slide_layout = prs.slide_layouts[6]  # Use the layout for a blank slide
    slide = prs.slides.add_slide(slide_layout)
    return slide


def add_slide_title(presentation, slide_index, title_text):
    slide = presentation.slides[slide_index]
    title_placeholder = slide.shapes.title
    title_placeholder.text = title_text


def add_text_to_slide(slide, text):
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    paragraph = textbox.text_frame.add_paragraph()
    paragraph.text = text


def add_image_to_slide(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(image_path, left, top, width, height)


def add_chart(slide, data, chart_type):
    chart = slide.shapes.add_chart(
        chart_type=chart_type,
        x=Inches(1), y=Inches(1), cx=Inches(8), cy=Inches(5)
    ).chart

    chart_data = chart.chart_data
    categories = data[0]
    series_data = data[1:]

    for category in categories:
        chart_data.categories.append(category)

    for series in series_data:
        series_values = series[1:]
        chart_data.add_series(series[0], tuple(series_values))

    return chart


def add_table_to_slide(slide, rows, columns):
    table_shape = slide.shapes.add_table(rows, columns)
    return table_shape.table


def add_shape(slide, left, top, width, height, shape_type):
    shape = slide.shapes.add_shape(
        shape_type=shape_type,
        left=Inches(left),
        top=Inches(top),
        width=Inches(width),
        height=Inches(height)
    )
    
    return shape


def add_hyperlink(slide, text, url):
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
    
    p = textbox.text_frame.add_paragraph()
    r = p.add_run()
    
    r.text = text
    r.hyperlink.address = url


def set_font_style_and_size(presentation, slide_number, font_style, font_size):
    slide = presentation.slides[slide_number - 1]

    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_style
                    run.font.size = Pt(font_size)


def set_slide_background_color(slide, color):
    r, g, b = color
    fill_format = slide.background.fill
    fill_format.solid()
    fill_format.fore_color.rgb = RGBColor(r,g,b)


def set_text_alignment(slide, text_box_index, align):
    
   
   
   # Get the text frame of the specified text box

   text_frame=slide.shapes[text_box_index].text_frame  # Set the alignment of the text
    
   if  align =='left':
       
      for paragraph in text_frame.paragraphs:           
          paragraph.alignment=PP_ALIGN.LEFT
      
   elif align=='center':
      
       for paragraph in text_frame.paragraphs:           
           paragraph.alignment=PP_ALIGN.CENTER
        
   elif align=='right':
       
       for paragraph in text_frame.paragraphs:           
           paragraph.alignment=PP_ALIGN.RIGHT
    
   else: 
       raise ValueError("Invalid alignment option. Choose from 'left', 'center', or 'right'.")
 


def set_border_style(slide_object,border_style):

        
       obj=slide_object
        
       obj.border_style=border_style

       
       return obj



def save_presentation_as_file(presentation ,filename ):
    
 presentation.save(filename)


def open_presentation(filepath ):

      prs=(filepath)
      
      return prs



def create_tree_visualization(presentation,data ,position,size):

     # Create a new slide
   
     Slide=presentation.slides.Add_slide(presentation.slide_layouts[1])# Assuming layout index 1 for blank slide
      
       
     #Add a shape to hold the tree visualization
    
     shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  position[0],position[1],size[0],size[1])
                                  

 
     #Format the shape as desired
    
     shape.fill.solid()# Use solid fill color
    
     shape.fill.fore_color.rgb=RGBColor(255 ,255 ,255)# Set fill color to white
      
       
     #Add text to represent the tree nodes
    
     

for node in data:

       text_frame=shape.text_frame.Add_paragraph()

       
       text_frame.text=node['text']
        
       
if 'children' in node:

            
for child_node in node['children']:

                 child_text_frame=text_frame.Add_paragraph()

                 
                 child_text_frame.text=child_node['text']

                 
if 'children' in child_node:

                       
for grandchild_node in child_node['children']:

                                grandchild_text_frame=
                                child_text_frame.Add_paragraph()


                                grandchild_text_frame.Text=
                                grandchild_node['text']


                                

 def add_node(presentation ,slide_number,node_text,parent_node=None):

     
      Slide=presentation.slides[slide_number-1]

      
      new_shape=
                   Slide.Shapes.Add_shape(MSO_SHAPE.RECTANGLE,
                                          left-100,
                                          top-100,width-150,height-50)

     
     
      new_shape.Text=node_text
      
     
if parent_node:
                                        
connector=
         Slide.Shapes.Add_connector(MSO_CONNECTOR.STRAIGHT,

         begin_x-parent_node.left+parent_node.width,

         begin_y-parent_node.top+parent_node.height/2,


end_x-new_shape.left,end_y-new_shape.top+new_shape.height/2)



return new_shape



 def remove_nodes(tree,nodes):

      
for node in nodes:
        
if node in tree:

               del tree[node]


return tree



 def rearrange_nodes_ppt(tree_structure,new_order):

 
return updated_tree_structure



 def customize_tree(tree,node_styles ,branch_styles):

     
for node_id ,style in node_styles.items():
        
if node_idin tree.nodes:


tree.nodes[node_id].style.update(style)
     

for branch_id ,stylein branch_styles.items():
        

if branch_idin tree.branches:


tree.branches[branch_id].style.update(style)

     
return tree



 def create_org_chart(chart_data ):



prs=

 Presentation()


Slide_layout=

prs.Slide_layouts [5]

Slide=

prs.Slides.Add_Slide(Slide_layout)


x,y,cx ,cy-Inches (2),

            
Inches (1),

            
Inches (6),

            
Inches (4.5)

chart=

Slide.Shapes.Add_chart(type-prs.chart_type.CHART_TYPE_ORG,x-x,y-y,cx-cx ,cy-cy )

chart

chart_data_excel-chart_data.to_exce()# assuming 'chart_data' is a pandas DataFrame with the org chart data

worksheet=

chart.chart_data.workbook.sheets [0]

for row_indexin range(len(chart_data_excel )):

for col_indexin range(len(chart_data_excel.columns)):


cell_value-chart_data_excel.iloc[row_index,col_index ]

worksheet.cell(row-row_index+1,column-col_index+1).value-cell_value



node_count-len(chart_data)-

foriin range(node_count ):

node-chart.plots [0].Add_category()

node.label.Text_Frame.Text-f"{chart_data ['Name'][i+1]}"

node.label.Text_Frame.Paragraphs [0].Alignment-PP_ALIGN.CENTER



parent_node_id-int(chart_data ['ParentID'][i+1 ])

 parent_node_id>0:


connector-chart.plots [0].connect_nodes(i,parent_node_id -1)



prs.save("organizational_chart.pptx")

 
 
 
 
 
 def arrange_nodes(tree_structure ):



prs-Presentation()


Slide-prs.Slides.Add_Slide(prs.Slide_layouts [6])


left-top-Inches (4),

width-height-Inches (8)


shape-Slide.Shapes .Add_shape,

MSO_SHAPE.RECTANGLE,left ,top,width,height )



shape.Fill.solid(),

shape.Fill.Fore_color.rgb-RGBColor (255 ,255 ,255)

shape.line.color.rgb-RGBColor (0,

Apply predefined rules or algorithms to arrange and format nodes


prs.save("hierarchical_visualization.pptx")

 



 def export_as_image(presentation ,slide_number):

Slide-presentation.Slides [slide_number -]


Slide.export('visualization.png ','png')






 def import_csv_data(file_path ):



hierarchical_data-'[]'


with open(file_path ,'r')as csvfile:


reader-csv.reader(csvfile)



for rowin reader:



hierarchical_data.append(row)



return hierarchical_data





 def export_to_csv(data ,filename ):


df-pd.DataFrame(data)


df.to_csv(filename,index-False)

