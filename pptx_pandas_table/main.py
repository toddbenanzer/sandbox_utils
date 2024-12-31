from .exceptions import DataAnalyzerError
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_ALIGN_PARAGRAPH
from pptx.table import Table
from pptx.util import Inches
import logging
import os
import pandas as pd


class PowerPointIntegration:
    """
    A class to handle the integration of pandas DataFrames into PowerPoint presentations.
    """

    def __init__(self, presentation_file: str):
        """
        Initialize the PowerPointIntegration class with a specific PowerPoint presentation.

        Args:
            presentation_file (str): The path to the PowerPoint file to be edited or created.

        Raises:
            FileNotFoundError: If the presentation file cannot be found.
        """
        try:
            self.presentation = Presentation(presentation_file)
        except Exception as e:
            raise FileNotFoundError(f"Could not open presentation file: {presentation_file}. Error: {str(e)}")

    def add_dataframe_as_table(self, dataframe: pd.DataFrame, slide_number: int, **kwargs):
        """
        Add a pandas DataFrame as a formatted table to a specified slide in the PowerPoint presentation.

        Args:
            dataframe (pd.DataFrame): The pandas DataFrame to be inserted as a table.
            slide_number (int): The index of the slide where the table should be added.
            **kwargs: Additional options for customization.
                      Example options include `table_style` for table styling and `number_format` for custom number formatting.

        Raises:
            IndexError: If the slide number is out of range.
            ValueError: If the DataFrame is empty or invalid.
        """
        # Validate slide number
        if slide_number < 0 or slide_number >= len(self.presentation.slides):
            raise IndexError(f"Slide number {slide_number} is out of range. Presentation has {len(self.presentation.slides)} slides.")

        # Validate DataFrame
        if dataframe.empty:
            raise ValueError("The DataFrame is empty and cannot be added as a table.")
        
        # Access the slide
        slide = self.presentation.slides[slide_number]

        rows, cols = dataframe.shape
        top = Inches(2.0)  # Example placeholder
        left = Inches(2.0)  # Example placeholder
        width = Inches(6.0)  # Example placeholder
        height = Inches(4.0)  # Example placeholder

        table = slide.shapes.add_table(rows + 1, cols, top, left, width, height).table

        # Set column headings
        for i, column_name in enumerate(dataframe.columns):
            table.cell(0, i).text = column_name

        # Populate table cells with DataFrame values
        for i, (index, row) in enumerate(dataframe.iterrows()):
            for j, value in enumerate(row):
                table.cell(i + 1, j).text = str(value)

        # Apply additional formatting from kwargs if provided
        table_style = kwargs.get('table_style', None)
        if table_style:
            # Apply the specified table style (assumes TableFormatter exists)
            # table_formatter = TableFormatter(table, dataframe)
            # table_formatter.apply_style(table_style)
            pass
        
        number_format = kwargs.get('number_format', None)
        if number_format:
            # Apply the specified number formats
            # table_formatter.apply_number_formatting(number_format)
            pass



class TableFormatter:
    """
    A class to handle the formatting of PowerPoint tables, specifically for number formatting and applying styles.
    """

    def __init__(self, table: Table, dataframe: pd.DataFrame):
        """
        Initialize the TableFormatter class with a PowerPoint table and a corresponding pandas DataFrame.

        Args:
            table (pptx.table.Table): The PowerPoint table object to be formatted.
            dataframe (pandas.DataFrame): The pandas DataFrame that corresponds to the data in the table.
        """
        self.table = table
        self.dataframe = dataframe

    def apply_number_formatting(self, column_formats: dict):
        """
        Apply specified number formatting to columns within the table based on the DataFrame.

        Args:
            column_formats (dict): A dictionary specifying the desired number format for each column,
                                   e.g., {'Sales': '${:,.2f}'}.

        Raises:
            KeyError: If a specified column in column_formats doesn't exist in the DataFrame.
        """
        for column_name, format_spec in column_formats.items():
            if column_name in self.dataframe.columns:
                col_index = self.dataframe.columns.get_loc(column_name)
                for i, value in enumerate(self.dataframe[column_name]):
                    formatted_value = format_spec.format(value)
                    self.table.cell(i + 1, col_index).text = formatted_value
            else:
                raise KeyError(f"Column '{column_name}' does not exist in the DataFrame.")

    def apply_style(self, style_name: str):
        """
        Apply a visual style to the table.

        Args:
            style_name (str): The name of the style to be applied to the table (e.g., 'Light Style 1 - Accent 1').
        """
        # Example basic style application
        # Here we use a simple logic, real-life implementation might need the style from a library or a database

        # For demonstration purposes, let's assume style_name dictates font and fill color
        if style_name == 'Light Style 1 - Accent 1':
            fill_color = RGBColor(242, 242, 242)  # Light gray
            font_color = RGBColor(0, 0, 0)       # Black
        elif style_name == 'Dark Style 1 - Accent 2':
            fill_color = RGBColor(68, 114, 196)  # Dark blue
            font_color = RGBColor(255, 255, 255) # White
        else:
            fill_color = RGBColor(255, 255, 255) # Default white
            font_color = RGBColor(0, 0, 0)       # Default black

        # Apply styling to each cell
        for row in self.table.rows:
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_color

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in paragraph.runs:
                        run.font.color.rgb = font_color

                # Additionally set vertical alignment
                cell.vertical_anchor = WD_CELL_VERTICAL_ALIGNMENT.MIDDLE


class UserInputHandler:
    """
    A class to handle user input for obtaining formatting preferences for DataFrame columns when creating PowerPoint tables.
    """

    def __init__(self):
        """
        Initialize the UserInputHandler class.
        """
        self.preferences = {}

    def get_column_formatting_preferences(self):
        """
        Obtain user-specified formatting preferences for each column.

        Returns:
            dict: A dictionary where each key is a column name and each value is the
                  user's preferred formatting string (e.g., {'Revenue': '${:,.2f}'}).
        """
        print("Enter formatting preferences for the columns in the format '{:format_specification}'.")
        print("For example, use '${:,.2f}' for currency formatting with two decimal places.")
        print("Enter 'done' when finished.")

        while True:
            column_name = input("Enter column name (or 'done' to finish): ")
            if column_name.lower() == 'done':
                break
            
            format_spec = input(f"Enter format specification for '{column_name}': ")
            
            # Example simple check for format specification validity
            # A more sophisticated check could be added depending on requirements
            try:
                test_value = 1234.56  # sample test value for format verification
                format_spec.format(test_value)
                self.preferences[column_name] = format_spec
                print(f"Format for column '{column_name}': {format_spec} stored successfully.")
            except (ValueError, KeyError) as e:
                print(f"Invalid format specification. Error: {str(e)}. Please try again.")

        return self.preferences


def format_numbers(data, format_spec):
    """
    Format numerical data according to a specified format string.

    Args:
        data (list): A collection of numerical values to be formatted.
        format_spec (str): A format string that outlines how each numerical value within `data`
                           should be formatted (e.g., "{:,.2f}").

    Returns:
        list: A list of formatted strings where each string corresponds to a formatted value from
              `data` as per `format_spec`.

    Raises:
        ValueError: If the data contains non-numeric elements that cannot be formatted.
        TypeError: If the format_spec is not a valid format string.
    """
    formatted_data = []

    for value in data:
        try:
            # Attempt to format the value using the specified format_spec
            formatted_value = format_spec.format(value)
            formatted_data.append(formatted_value)
        except (ValueError, TypeError) as e:
            raise ValueError(f"Error formatting value '{value}' with the format '{format_spec}': {str(e)}")

    return formatted_data


def get_table_styles():
    """
    Retrieve a list of available styles that can be applied to PowerPoint tables.

    Returns:
        list of dict: A list of dictionaries, each representing a table style.
                      Each dictionary contains the 'name' and 'description' of the style.
    """
    # Example styles; these should match those available from PowerPoint or
    # be customizable styles that your application supports
    styles = [
        {
            "name": "Light Style 1 - Accent 1",
            "description": "Light style with a subtle accent on the first column."
        },
        {
            "name": "Medium Style 2 - Accent 2",
            "description": "Medium style with stripe rows and emphasis on the second column."
        },
        {
            "name": "Dark Style 3 - Accent 3",
            "description": "Dark style designed to stand out with a strong accent."
        },
        {
            "name": "Grid Table 4 - Accent 4",
            "description": "Grid style with all-around borders for clear separation."
        }
    ]

    # The actual implementation might fetch this information from a file or resources
    # associated with PowerPoint or user-defined style settings

    return styles



def validate_dataframe(dataframe: pd.DataFrame) -> bool:
    """
    Validate the provided pandas DataFrame to ensure it meets the necessary criteria for further operations.

    Args:
        dataframe (pandas.DataFrame): The DataFrame to be validated.

    Returns:
        bool: True if the DataFrame is valid, otherwise raises an appropriate exception.

    Raises:
        ValueError: If the DataFrame is empty, has zero columns, or has duplicate/empty column names.
        TypeError: If the DataFrame contains unsupported data types.
    """
    # Check if the DataFrame is empty
    if dataframe.empty:
        raise ValueError("The DataFrame is empty. It must contain at least one row and one column.")

    # Check for a valid number of columns
    if dataframe.shape[1] == 0:
        raise ValueError("The DataFrame contains no columns.")

    # Check for unique and non-empty column names
    if not all(dataframe.columns):
        raise ValueError("Some columns have empty names. Column names must be non-empty.")
    if len(dataframe.columns) != len(set(dataframe.columns)):
        raise ValueError("Some columns have duplicate names. Column names must be unique.")

    # Check for unsupported data types
    supported_types = {int, float, bool, str}
    for dtype in dataframe.dtypes:
        if dtype.type not in supported_types:
            raise TypeError(f"The DataFrame contains unsupported data types: {dtype}")
    
    return True



def export_presentation(presentation: Presentation, file_path: str) -> None:
    """
    Save a PowerPoint presentation object to a specified file path.

    Args:
        presentation (pptx.Presentation): The PowerPoint presentation object to be saved.
        file_path (str): The file path where the presentation should be saved. Must include the .pptx extension.

    Raises:
        FileNotFoundError: If the specified directory in the file_path cannot be found.
        PermissionError: If there are insufficient permissions to write to the file_path.
        ValueError: If the file_path does not end with .pptx.
    """
    # Validate the file extension
    if not file_path.endswith('.pptx'):
        raise ValueError("The file path must end with a '.pptx' extension.")
    
    # Ensure the directory exists
    directory = os.path.dirname(file_path)
    if not os.path.exists(directory):
        try:
            os.makedirs(directory)
        except Exception as e:
            raise FileNotFoundError(f"Cannot create directory {directory}. Error: {str(e)}")
    
    # Attempt to save the presentation
    try:
        presentation.save(file_path)
    except PermissionError:
        raise PermissionError(f"Insufficient permissions to write to {file_path}.")
    except Exception as e:
        raise Exception(f"An error occurred while saving the presentation: {str(e)}")



def setup_logger(log_level: int) -> logging.Logger:
    """
    Configure the logger for the application.

    Args:
        log_level (int): Defines the level of logging detail to capture (e.g., logging.DEBUG)

    Returns:
        logging.Logger: A configured logger instance.
    """
    logger = logging.getLogger("pptx_formatter")
    logger.setLevel(log_level)

    # Create console handler
    console_handler = logging.StreamHandler()
    console_handler.setLevel(log_level)

    # Create formatter
    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')

    # Add formatter to console handler
    console_handler.setFormatter(formatter)

    # Add handler to logger
    logger.addHandler(console_handler)

    # Optionally add file handler
    # file_handler = logging.FileHandler('app.log')
    # file_handler.setFormatter(formatter)
    # logger.addHandler(file_handler)

    return logger
