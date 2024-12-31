from mymodule import PowerPointIntegration  # Replace 'mymodule' with the actual module name
from mymodule import TableFormatter  # Replace 'mymodule' with the actual module name
from mymodule import UserInputHandler  # Replace 'mymodule' with the actual module name
from mymodule import export_presentation  # Replace 'mymodule' with the actual module name
from mymodule import setup_logger  # Replace 'mymodule' with the actual module name
from mymodule import validate_dataframe  # Replace 'mymodule' with the actual module name
from pptx import Presentation
from pptx.util import Inches
import logging
import os
import pandas as pd


# Example 1: Adding a simple DataFrame as a table to a new slide
dataframe = pd.DataFrame({
    "Name": ["Alice", "Bob"],
    "Age": [25, 30],
    "Salary": [50000, 60000]
})

# Create a new presentation
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Use a slide layout with a title and a content
presentation.slides.add_slide(slide_layout)
presentation_file = "presentation1.pptx"
presentation.save(presentation_file)

# Integrate the DataFrame into the presentation
ppt_integration = PowerPointIntegration(presentation_file)
ppt_integration.add_dataframe_as_table(dataframe, slide_number=0)

# Save the presentation
ppt_integration.presentation.save("output_presentation1.pptx")

# Example 2: Adding a DataFrame with custom formatting to a specified slide
dataframe2 = pd.DataFrame({
    "Product": ["Widget", "Gadget"],
    "Units Sold": [100, 150],
    "Revenue": [2000.5, 3000.75]
})

# Create another presentation
presentation2 = Presentation()
presentation2.slides.add_slide(slide_layout)
presentation_file2 = "presentation2.pptx"
presentation2.save(presentation_file2)

# Integrate the DataFrame with custom formatting
ppt_integration2 = PowerPointIntegration(presentation_file2)
ppt_integration2.add_dataframe_as_table(
    dataframe2,
    slide_number=0,
    table_style="Light Style 1 - Accent 1",
    number_format={"Revenue": "${:,.2f}"}
)

# Save the presentation
ppt_integration2.presentation.save("output_presentation2.pptx")



# Example 1: Simple number formatting for a Sales table
dataframe = pd.DataFrame({
    "Product": ["Widget", "Gadget"],
    "Sales": [2000.5, 3000.75]
})

presentation = Presentation()
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
table = slide.shapes.add_table(3, 2, Inches(2.0), Inches(2.0), Inches(6.0), Inches(4.0)).table

# Insert header row
table.cell(0, 0).text = "Product"
table.cell(0, 1).text = "Sales"

# Insert data rows
for i, row in dataframe.iterrows():
    table.cell(i + 1, 0).text = str(row["Product"])
    table.cell(i + 1, 1).text = str(row["Sales"])

table_formatter = TableFormatter(table, dataframe)

# Apply number formatting
table_formatter.apply_number_formatting({'Sales': '${:,.2f}'})
presentation.save('presentation_with_number_formatting.pptx')

# Example 2: Applying a light style to the table
table_formatter.apply_style('Light Style 1 - Accent 1')
presentation.save('presentation_with_light_style.pptx')

# Example 3: Applying a dark style to the table
table_formatter.apply_style('Dark Style 1 - Accent 2')
presentation.save('presentation_with_dark_style.pptx')



# Example 1: Basic use with hypothetical user inputs via a test setup
# Simulate user input: ["Revenue", "${:,.2f}", "Cost", "{:.1f}", "done"]

user_handler = UserInputHandler()
preferences = user_handler.get_column_formatting_preferences()
print("User Preferences:", preferences)  # Expected output: {'Revenue': '${:,.2f}', 'Cost': '{:.1f}'}

# Example 2: Handling no formatting input (user directly finishes)
# Simulate user input: ["done"]

user_handler_empty = UserInputHandler()
empty_preferences = user_handler_empty.get_column_formatting_preferences()
print("User Preferences when done immediately:", empty_preferences)  # Expected output: {}

# Example 3: Demonstrating invalid input correction
# Simulate user input: ["Invalid", "invalid_format", "{:,.2f}", "done"]

user_handler_invalid = UserInputHandler()
preferences_with_correction = user_handler_invalid.get_column_formatting_preferences()
print("Preferences with corrected inputs:", preferences_with_correction)  # Expected output: {'Invalid': '{:,.2f}'}


# Example 1: Basic usage with commas and two decimal places
data1 = [1234.56, 7890, 12345.678]
format_spec1 = "{:,.2f}"
formatted_numbers1 = format_numbers(data1, format_spec1)
print(formatted_numbers1)  # Output: ['1,234.56', '7,890.00', '12,345.68']

# Example 2: Using no decimal places
data2 = [1234.56, 7890, 12345.678]
format_spec2 = "{:,.0f}"
formatted_numbers2 = format_numbers(data2, format_spec2)
print(formatted_numbers2)  # Output: ['1,235', '7,890', '12,346']

# Example 3: Percentage format
data3 = [0.1, 0.23, 0.456]
format_spec3 = "{:.1%}"
formatted_numbers3 = format_numbers(data3, format_spec3)
print(formatted_numbers3)  # Output: ['10.0%', '23.0%', '45.6%']

# Example 4: Handling different number styles, e.g., scientific notation
data4 = [1234.56, 7890, 12345.678]
format_spec4 = "{:.2e}"
formatted_numbers4 = format_numbers(data4, format_spec4)
print(formatted_numbers4)  # Output: ['1.23e+03', '7.89e+03', '1.23e+04']

# Example 5: Custom currency format with leading text
data5 = [1234.56, 7890, 12345.678]
format_spec5 = "${:,.2f} USD"
formatted_numbers5 = format_numbers(data5, format_spec5)
print(formatted_numbers5)  # Output: ['$1,234.56 USD', '$7,890.00 USD', '$12,345.68 USD']


# Example 1: Retrieve and print all table styles
styles = get_table_styles()
for style in styles:
    print(f"Style Name: {style['name']}, Description: {style['description']}")

# Example 2: Check for the existence of a specific style
def is_style_available(style_name):
    return any(style['name'] == style_name for style in get_table_styles())

style_to_check = "Dark Style 3 - Accent 3"
print(f"Is '{style_to_check}' available? {is_style_available(style_to_check)}")

# Example 3: Get descriptions of all available styles
def get_style_descriptions():
    return [style['description'] for style in get_table_styles()]

descriptions = get_style_descriptions()
print("Descriptions of available styles:", descriptions)



# Example 1: Valid DataFrame
df_valid = pd.DataFrame({
    "A": [1, 2, 3],
    "B": [4.5, 5.5, 6.5],
    "C": ["foo", "bar", "baz"]
})

try:
    result = validate_dataframe(df_valid)
    print("Example 1: DataFrame is valid? ->", result)  # Output: True
except (ValueError, TypeError) as e:
    print("Example 1: Validation failed with message:", e)

# Example 2: DataFrame with Duplicate Column Names
df_duplicate_columns = pd.DataFrame({
    "A": [1, 2, 3],
    "A": [4, 5, 6]
})

try:
    validate_dataframe(df_duplicate_columns)
except (ValueError, TypeError) as e:
    print("Example 2: Validation failed with message:", e)  # Expect an error about duplicate column names

# Example 3: DataFrame with Unsupported Data Type
df_unsupported_dtype = pd.DataFrame({
    "A": [1, 2, 3],
    "B": [pd.Timestamp('2021-01-01'), pd.Timestamp('2021-01-02'), pd.Timestamp('2021-01-03')]
})

try:
    validate_dataframe(df_unsupported_dtype)
except (ValueError, TypeError) as e:
    print("Example 3: Validation failed with message:", e)  # Expect an error about unsupported data types

# Example 4: Empty DataFrame
df_empty = pd.DataFrame()

try:
    validate_dataframe(df_empty)
except (ValueError, TypeError) as e:
    print("Example 4: Validation failed with message:", e)  # Expect an error about the DataFrame being empty



# Example 1: Save a simple empty presentation to a valid path
presentation = Presentation()
slide_layout = presentation.slide_layouts[0]  # Choose a layout from PowerPoint
slide = presentation.slides.add_slide(slide_layout)  # Add one slide
file_path = os.path.join(os.getcwd(), "presentation1.pptx")
export_presentation(presentation, file_path)

# Example 2: Save a presentation with a new slide and text to a nested directory
presentation2 = Presentation()
slide_layout2 = presentation2.slide_layouts[0]
slide2 = presentation2.slides.add_slide(slide_layout2)
title = slide2.shapes.title
title.text = "Hello, World!"
file_path2 = os.path.join(os.getcwd(), "nested/directory/presentation2.pptx")
export_presentation(presentation2, file_path2)

# Example 3: Attempt to save with an invalid file extension, demonstrating error handling
try:
    export_presentation(presentation2, "presentation3.txt")
except ValueError as e:
    print("Example 3 Error:", e)

# Example 4: Ensure directory creation and presentation saving in a specific folder
specific_folder_path = os.path.join(os.getcwd(), "specific_folder")
file_path4 = os.path.join(specific_folder_path, "presentation4.pptx")
export_presentation(presentation, file_path4)



# Example 1: Setup logger with DEBUG level
logger_debug = setup_logger(logging.DEBUG)
logger_debug.debug("This is a debug message with DEBUG level.")
logger_debug.info("This message is also visible since it's INFO level.")

# Example 2: Setup logger with INFO level
logger_info = setup_logger(logging.INFO)
logger_info.info("This is an informational message with INFO level.")
logger_info.debug("This debug message will not be visible.")

# Example 3: Setup logger with WARNING level
logger_warning = setup_logger(logging.WARNING)
logger_warning.warning("This is a warning message with WARNING level.")
logger_warning.info("This info message will not be visible.")

# Example 4: Setup logger with ERROR level
logger_error = setup_logger(logging.ERROR)
logger_error.error("This is an error message with ERROR level.")
logger_error.warning("This warning will not be visible.")
