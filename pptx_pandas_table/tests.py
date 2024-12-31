from mymodule import PowerPointIntegration  # Replace 'mymodule' with the actual module name
from mymodule import TableFormatter  # Replace 'mymodule' with the actual module name
from mymodule import UserInputHandler  # Replace 'mymodule' with the actual module name
from mymodule import export_presentation  # Replace 'mymodule' with the actual module name
from mymodule import format_numbers  # Replace 'mymodule' with the actual module name
from mymodule import get_table_styles  # Replace 'mymodule' with the actual module name
from mymodule import setup_logger  # Replace 'mymodule' with the actual module name
from mymodule import validate_dataframe  # Replace 'mymodule' with the actual module name
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from tempfile import TemporaryDirectory
from unittest.mock import patch, call
import logging
import os
import pandas as pd
import pytest


def test_init_valid_presentation():
    with TemporaryDirectory() as tmpdir:
        temp_pptx = Path(tmpdir) / "test_presentation.pptx"
        presentation = Presentation()
        presentation.save(temp_pptx)

        ppt_integration = PowerPointIntegration(str(temp_pptx))
        assert isinstance(ppt_integration.presentation, Presentation)

def test_init_invalid_presentation():
    with pytest.raises(FileNotFoundError):
        PowerPointIntegration("non_existent.pptx")

def test_add_dataframe_as_table_valid():
    dataframe = pd.DataFrame({
        "A": [1, 2],
        "B": [3, 4]
    })
    
    with TemporaryDirectory() as tmpdir:
        temp_pptx = Path(tmpdir) / "test_presentation.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[5])
        presentation.save(temp_pptx)

        ppt_integration = PowerPointIntegration(str(temp_pptx))
        ppt_integration.add_dataframe_as_table(dataframe, 0)

        slide = ppt_integration.presentation.slides[0]
        assert len(slide.shapes) > 0
        table = slide.shapes[0].table

        # Check header
        assert table.cell(0, 0).text == "A"
        assert table.cell(0, 1).text == "B"

        # Check data
        assert table.cell(1, 0).text == "1"
        assert table.cell(1, 1).text == "3"
        assert table.cell(2, 0).text == "2"
        assert table.cell(2, 1).text == "4"

def test_add_dataframe_as_table_invalid_slide_number():
    dataframe = pd.DataFrame({
        "A": [1]
    })

    with TemporaryDirectory() as tmpdir:
        temp_pptx = Path(tmpdir) / "test_presentation.pptx"
        presentation = Presentation()
        presentation.save(temp_pptx)

        ppt_integration = PowerPointIntegration(str(temp_pptx))
        with pytest.raises(IndexError):
            ppt_integration.add_dataframe_as_table(dataframe, 1)

def test_add_dataframe_as_table_empty_dataframe():
    dataframe = pd.DataFrame()

    with TemporaryDirectory() as tmpdir:
        temp_pptx = Path(tmpdir) / "test_presentation.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[5])
        presentation.save(temp_pptx)

        ppt_integration = PowerPointIntegration(str(temp_pptx))
        with pytest.raises(ValueError):
            ppt_integration.add_dataframe_as_table(dataframe, 0)



@pytest.fixture
def sample_data():
    dataframe = pd.DataFrame({
        "Product": ["Widget", "Gadget"],
        "Sales": [2000.5, 3000.75]
    })

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    rows, cols = dataframe.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(2.0), Inches(2.0), Inches(6.0), Inches(4.0)).table

    for i, column_name in enumerate(dataframe.columns):
        table.cell(0, i).text = column_name

    for i, (index, row) in enumerate(dataframe.iterrows()):
        for j, value in enumerate(row):
            table.cell(i + 1, j).text = str(value)

    return table, dataframe

def test_apply_number_formatting(sample_data):
    table, dataframe = sample_data
    formatter = TableFormatter(table, dataframe)
    
    column_formats = {'Sales': '${:,.2f}'}
    formatter.apply_number_formatting(column_formats)
    
    assert table.cell(1, 1).text == "$2,000.50"
    assert table.cell(2, 1).text == "$3,000.75"

def test_apply_number_formatting_invalid_column(sample_data):
    table, dataframe = sample_data
    formatter = TableFormatter(table, dataframe)
    
    column_formats = {'InvalidColumn': '${:,.2f}'}
    with pytest.raises(KeyError):
        formatter.apply_number_formatting(column_formats)

def test_apply_style_light(sample_data):
    table, dataframe = sample_data
    formatter = TableFormatter(table, dataframe)

    formatter.apply_style('Light Style 1 - Accent 1')
    for row in table.rows:
        for cell in row.cells:
            assert cell.fill.fore_color.rgb == RGBColor(242, 242, 242)

def test_apply_style_dark(sample_data):
    table, dataframe = sample_data
    formatter = TableFormatter(table, dataframe)

    formatter.apply_style('Dark Style 1 - Accent 2')
    for row in table.rows:
        for cell in row.cells:
            assert cell.fill.fore_color.rgb == RGBColor(68, 114, 196)

def test_apply_style_default(sample_data):
    table, dataframe = sample_data
    formatter = TableFormatter(table, dataframe)

    formatter.apply_style('Nonexistent Style')
    for row in table.rows:
        for cell in row.cells:
            assert cell.fill.fore_color.rgb == RGBColor(255, 255, 255)




def test_get_column_formatting_preferences_valid_input():
    user_input = [
        "Revenue", "${:,.2f}",
        "Cost", "{:.1f}",
        "done"
    ]
    
    expected_output = {
        "Revenue": "${:,.2f}",
        "Cost": "{:.1f}"
    }
    
    with patch('builtins.input', side_effect=user_input):
        handler = UserInputHandler()
        result = handler.get_column_formatting_preferences()
    
    assert result == expected_output

def test_get_column_formatting_preferences_invalid_format():
    user_input = [
        "Revenue", "invalid_format",  # Invalid format
        "${:,.2f}",  # Corrected valid format
        "done"
    ]
    
    expected_output = {
        "Revenue": "${:,.2f}"
    }
    
    with patch('builtins.input', side_effect=user_input), \
         patch('builtins.print') as mock_print:
        handler = UserInputHandler()
        result = handler.get_column_formatting_preferences()
    
    # Ensure the invalid format error is printed
    print_calls = [
        call("Invalid format specification. Error: 'invalid_format'. Please try again.")
    ]
    
    mock_print.assert_has_calls(print_calls, any_order=True)
    assert result == expected_output

def test_get_column_formatting_preferences_complete_without_errors():
    user_input = [
        "done"
    ]
    
    expected_output = {}
    
    with patch('builtins.input', side_effect=user_input), \
         patch('builtins.print') as mock_print:
        handler = UserInputHandler()
        result = handler.get_column_formatting_preferences()

    assert result == expected_output



def test_format_numbers_valid_data():
    data = [1000, 2000.5, 3000.75]
    format_spec = "{:,.2f}"
    expected = ["1,000.00", "2,000.50", "3,000.75"]
    assert format_numbers(data, format_spec) == expected

def test_format_numbers_different_format():
    data = [1000, 2000.5, 3000.75]
    format_spec = "{:.1f}"
    expected = ["1000.0", "2000.5", "3000.8"]
    assert format_numbers(data, format_spec) == expected

def test_format_numbers_empty_list():
    data = []
    format_spec = "{:.2f}"
    expected = []
    assert format_numbers(data, format_spec) == expected

def test_format_numbers_invalid_data_type():
    data = ["1000", 2000.5, "text"]
    format_spec = "{:,.2f}"
    with pytest.raises(ValueError) as excinfo:
        format_numbers(data, format_spec)
    assert "Error formatting value" in str(excinfo.value)

def test_format_numbers_invalid_format_spec():
    data = [1000, 2000.5, 3000.75]
    format_spec = "{:,.2a}"
    with pytest.raises(ValueError) as excinfo:
        format_numbers(data, format_spec)
    assert "Error formatting value" in str(excinfo.value)



def test_get_table_styles_returns_list():
    result = get_table_styles()
    assert isinstance(result, list), "Expected result to be a list."

def test_get_table_styles_non_empty():
    result = get_table_styles()
    assert len(result) > 0, "Expected result to have at least one style."

def test_get_table_styles_elements_are_dicts():
    result = get_table_styles()
    for style in result:
        assert isinstance(style, dict), "Expected each style to be a dictionary."

def test_get_table_styles_contains_correct_keys():
    result = get_table_styles()
    required_keys = {'name', 'description'}
    for style in result:
        assert required_keys.issubset(style.keys()), "Each style should contain 'name' and 'description' keys."

def test_get_table_styles_valid_names():
    result = get_table_styles()
    valid_names = {
        "Light Style 1 - Accent 1",
        "Medium Style 2 - Accent 2",
        "Dark Style 3 - Accent 3",
        "Grid Table 4 - Accent 4"
    }
    style_names = {style['name'] for style in result}
    assert valid_names == style_names, "Style names do not match the expected names."



def test_validate_dataframe_valid():
    df = pd.DataFrame({
        "A": [1, 2, 3],
        "B": [4.5, 5.5, 6.5],
        "C": ["x", "y", "z"]
    })
    assert validate_dataframe(df) is True

def test_validate_dataframe_empty():
    df = pd.DataFrame()
    with pytest.raises(ValueError, match=r"The DataFrame is empty"):
        validate_dataframe(df)

def test_validate_dataframe_no_columns():
    df = pd.DataFrame([], columns=[])
    with pytest.raises(ValueError, match=r"The DataFrame contains no columns"):
        validate_dataframe(df)

def test_validate_dataframe_empty_column_names():
    df = pd.DataFrame({
        "": [1, 2, 3],
        "B": [4.5, 5.5, 6.5]
    })
    with pytest.raises(ValueError, match=r"Some columns have empty names"):
        validate_dataframe(df)

def test_validate_dataframe_duplicate_column_names():
    df = pd.DataFrame({
        "A": [1, 2, 3],
        "A": [4, 5, 6]
    })
    with pytest.raises(ValueError, match=r"Some columns have duplicate names"):
        validate_dataframe(df)

def test_validate_dataframe_unsupported_data_type():
    df = pd.DataFrame({
        "A": [1, 2, 3],
        "B": [4.5, 5.5, 6.5],
        "D": [pd.Timestamp('20210101'), pd.Timestamp('20210102'), pd.Timestamp('20210103')]
    })
    with pytest.raises(TypeError, match=r"The DataFrame contains unsupported data types"):
        validate_dataframe(df)



@pytest.fixture
def dummy_presentation():
    # Create a simple presentation fixture
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[5])
    return presentation

def test_export_presentation_valid_path(tmp_path, dummy_presentation):
    file_path = tmp_path / "test_presentation.pptx"
    export_presentation(dummy_presentation, str(file_path))
    assert file_path.exists()

def test_export_presentation_invalid_extension(dummy_presentation):
    with pytest.raises(ValueError, match=r"The file path must end with a '.pptx' extension."):
        export_presentation(dummy_presentation, "invalid_path/test_presentation.txt")

def test_export_presentation_create_directories(tmp_path, dummy_presentation):
    new_dir = tmp_path / "new_folder"
    file_path = new_dir / "test_presentation.pptx"
    export_presentation(dummy_presentation, str(file_path))
    assert file_path.exists()
    assert os.path.exists(new_dir)

def test_export_presentation_permission_error(monkeypatch, tmp_path, dummy_presentation):
    # Mock an environment where the directory is not writable
    def raise_permission_error(*args, **kwargs):
        raise PermissionError("Mocked permission error")

    monkeypatch.setattr(Presentation, "save", raise_permission_error)

    file_path = tmp_path / "test_presentation.pptx"
    with pytest.raises(PermissionError, match=r"Insufficient permissions to write to"):
        export_presentation(dummy_presentation, str(file_path))



def test_setup_logger_debug_level(caplog):
    logger = setup_logger(logging.DEBUG)
    with caplog.at_level(logging.DEBUG, logger="pptx_formatter"):
        logger.debug("Debug message")
    assert "Debug message" in caplog.text

def test_setup_logger_info_level(caplog):
    logger = setup_logger(logging.INFO)
    with caplog.at_level(logging.INFO, logger="pptx_formatter"):
        logger.info("Info message")
        logger.debug("Debug message")
    assert "Info message" in caplog.text
    assert "Debug message" not in caplog.text

def test_setup_logger_warning_level(caplog):
    logger = setup_logger(logging.WARNING)
    with caplog.at_level(logging.WARNING, logger="pptx_formatter"):
        logger.warning("Warning message")
        logger.info("Info message")
    assert "Warning message" in caplog.text
    assert "Info message" not in caplog.text

def test_setup_logger_error_level(caplog):
    logger = setup_logger(logging.ERROR)
    with caplog.at_level(logging.ERROR, logger="pptx_formatter"):
        logger.error("Error message")
        logger.warning("Warning message")
    assert "Error message" in caplog.text
    assert "Warning message" not in caplog.text

def test_logger_has_console_handler():
    logger = setup_logger(logging.INFO)
    handlers = logger.handlers
    assert any(isinstance(handler, logging.StreamHandler) for handler in handlers)
