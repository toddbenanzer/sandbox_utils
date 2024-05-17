
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.table import MSO_TABLE_STYLE

def convert_dataframe_to_ppt_table(df):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    num_rows, num_cols = df.shape
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.8 * num_rows)
    table = slide.shapes.add_table(num_rows + 1, num_cols, left, top, width, height).table

    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = col_name
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(192, 192, 192)

    for r in range(num_rows):
        for c in range(num_cols):
            cell = table.cell(r + 1, c)
            value = df.iloc[r, c]
            cell.text = str(value)

            if pd.api.types.is_numeric_dtype(df.dtypes[c]):
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                cell.text_frame.clear()
                cell.text_frame.add_paragraph().text = '{:,.2f}'.format(value)

    return prs

def add_table_to_slide(prs_path, slide_index, df):
    prs = Presentation(prs_path)
    slide = prs.slides[slide_index]

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    left = Inches(1)
    top = Inches(1)
    width = slide_width - Inches(2)
    height = slide_height - Inches(2)

    table_shape = slide.shapes.add_table(rows=df.shape[0] + 1, cols=df.shape[1], left=left, top=top,
                                         width=width, height=height).table

    col_width = width / df.shape[1]
    
    for i, col in enumerate(df.columns):
        table_shape.columns[i].width = col_width
        header_cell(table_shape.cell(0, i), col)

        for j in range(df.shape[0]):
            data_cell(table_shape.cell(j + 1, i), df.iloc[j, i])

    prs.save("updated_presentation.pptx")

def header_cell(cell, text):
    cell.text_frame.clear()
    p = cell.text_frame.add_paragraph()
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font_size(p.runs[0], 14)

def data_cell(cell, value):
    cell.text_frame.clear()
    
    if isinstance(value, (int,float)):
        p_text_format(cell.text_frame.add_paragraph(), '{:,.2f}'.format(value), PP_ALIGN.RIGHT)    
        
def set_font_size(run,font_size):  
   run.font.size=Pt(font_size)    

def set_table_style(table, style):
   table.style=style 

def set_table_font_color(table , color): 
   for row in table.rows:
       for cell in row.cells:
          for paragraph in cell.text_frame.paragraphs:
              for run in paragraph.runs:
                   run.font.color.rgb=RGBColor.from_string(color)  

def set_table_background_color(table,color):
   for row in table.rows:
       for cell in row.cells:
          cell.fill.solid()
          cell.fill.fore_color.rgb=RGBColor.from_string(color)

def set_table_border_color(table,color): 
   for row in table.rows:
       for cell in row.cells:   
           for side in ['left','top','bottom' ,'right']:
               getattr(cell.border ,side).color.rgb=color 

def set_column_widths(table,widths):  
     assert len(widths)==len(table.columns),"column length must match the given widths"
     [setattr(column,'width',Inches(width))for column,width in zip (table.columns,widths)]

def set_table_row_heights(table,row_height): 
   [setattr(row,'height',Pt(row_height))for row in table.rows]

def add_currency_table(presentation ,df):  
   slide=presentation.slides.add_slide(presentation.slide_layouts[1])
   left=Inches(1)
   top=Inches (2)
   width=Inches (6)  
   height=Inches (4)  
   
   table_df_to_ppt(df ,slide,left ,top ,width,height,currency=True)

   return presentation  

#helper function to add DataFrame to Table 
def table_df_to_ppt(df ,slide,left ,top,width,height,currency=False): 

     table_shape=slide.shapes.add_table(rows=df.shape [0]+1 ,cols=df.shape [1],left=left,
                                             top=top,width=width,height=height).table
    
     [header_cell(table_shape.cell(0,i),col_name)for i,col_name in enumerate (df.columns)] 
   
     if currency :
         df=df.applymap(lambda x:f'${:,.2f}'.format(x)if isinstance(x,(float,int))else x )
     
     [[data_cell(table_shape.cell(i+1,j),df.iloc[i,j])for j in range(df.shape[ 1])]for i in range (df.shape[ 0])]


presentation.save('currency_table.pptx')

def format_column_as_percentage(table,column_index): 
  column_cells=[row.cells[column_index]for row in table.rows]
  [[p_text_format(cell .text_frame.paragraphs [0],'{:.2%}'.format(float(cell .text.strip('%'))/100),
                    PP_PARAGRAPH_ALIGNMENT.RIGHT)]for i ,cell in enumerate(column_cells [1:],start= 1)]
  
  
#helper function to add text and formatting to the paragraph.
def p_text_format(paragraph,text ,alignment=None,font_color=None,font_bold=False,font_size=None):  
  paragraph.clear()
  run=paragraph.add_run()   
  run.text=text
   
 if alignment:
      paragraph.alignment=getattr(PP_ALIGN if alignment.upper()[-3:]=='IGN'else PP_PARAGRAPH_ALIGNMENT,
                                  alignment.upper())
     
 if font_bold:
      run.font.bold=True
      
 if font_size:
      run.font.size=Pt(font_size)

presentation.save("formatted_table.pptx") 


