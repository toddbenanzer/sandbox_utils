from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches
from pptx.util import Pt
import logging
import os
import pandas as pd


class PresentationGenerator:
    """
    A class to generate PowerPoint presentations with customizable slides.
    """

    def __init__(self, title, **kwargs):
        """
        Initialize the PresentationGenerator with a title and optional settings.

        Args:
            title (str): The title of the presentation.
            **kwargs: Optional settings for the presentation.

        Raises:
            ValueError: If the title is not provided as a string.
        """
        if not isinstance(title, str):
            raise ValueError("Title must be a string.")
        
        self.presentation = Presentation()
        self.title = title
        self.slides = []
        
        # Optional: Apply template settings based on kwargs
        template = kwargs.get('template')
        if template:
            # TODO: Load template (this is a placeholder)
            pass

        # Add title slide
        title_slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = self.title
        self.slides.append(slide)

    def add_slide(self, title, content_type, **kwargs):
        """
        Add a slide to the presentation with specified content.

        Args:
            title (str): Title for the slide.
            content_type (str): Type of content ('table', 'text', etc.).
            **kwargs: Additional customization options for the slide.

        Raises:
            ValueError: If content_type is not valid.
        """
        if content_type not in ['table', 'text']:
            raise ValueError("Unsupported content_type. Choose 'table' or 'text'.")

        layout = self.presentation.slide_layouts[5]  # Using a blank slide layout
        slide = self.presentation.slides.add_slide(layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        if content_type == 'table':
            # TODO: Implement table creation logic
            pass
        elif content_type == 'text':
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5.5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_content = kwargs.get('text_content', '')
            text_frame.text = text_content

        self.slides.append(slide)

    def save_presentation(self, file_path):
        """
        Save the presentation to the specified file path.

        Args:
            file_path (str): The file path to save the presentation.

        Raises:
            IOError: If the file cannot be written.
            ValueError: If the file path is invalid.
        """
        if not isinstance(file_path, str) or not file_path.endswith('.pptx'):
            raise ValueError("Invalid file path. Must be a .pptx file.")
        try:
            self.presentation.save(file_path)
        except IOError as e:
            raise IOError(f"Failed to save presentation: {e}")



class DataFrameToSlide:
    """
    A class to convert a pandas DataFrame into a styled table within a PowerPoint slide.
    """

    def __init__(self, dataframe, style_config):
        """
        Initialize the DataFrameToSlide with the data and style configuration.

        Args:
            dataframe (pd.DataFrame): The data to be displayed in the slide.
            style_config (StyleConfig): Styling options for the table display.

        Raises:
            ValueError: If the dataframe is not a valid pandas DataFrame.
        """
        if not isinstance(dataframe, pd.DataFrame):
            raise ValueError("Invalid dataframe. Must be a pandas DataFrame.")
        
        self.dataframe = dataframe
        self.style_config = style_config

    def create_table_slide(self, slide):
        """
        Create a table slide in the specified PowerPoint slide.

        Args:
            slide (Slide): The PowerPoint slide to add the table to.

        Raises:
            ValueError: If the slide is not valid for accepting a table.
        """
        if not isinstance(slide, Slide):
            raise ValueError("Invalid slide provided.")

        rows, cols = self.dataframe.shape
        left = Inches(2.0)
        top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)
        
        table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

        # Setting up the header row
        for col, column_name in enumerate(self.dataframe.columns):
            table.cell(0, col).text = str(column_name)
        
        # Filling the table with data
        for row in range(rows):
            for col in range(cols):
                table.cell(row+1, col).text = str(self.dataframe.iloc[row, col])
        
        # Apply styles
        self.apply_styles(table, self.style_config)

    def apply_styles(self, table, style_config):
        """
        Apply the specified styles from the StyleConfig to the table.

        Args:
            table: The table object within the slide to which styles are applied.
            style_config (StyleConfig): The styling attributes to be applied to the table.
        """
        for row in range(len(table.rows)):
            for col in range(len(table.columns)):
                cell = table.cell(row, col)
                # Apply font style
                cell.text_frame.paragraphs[0].font.name = style_config.font_style
                # Apply font size
                cell.text_frame.paragraphs[0].font.size = style_config.font_size
                # Apply font color
                cell.text_frame.paragraphs[0].font.color.rgb = style_config.font_color
                # Align text
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT if row == 0 else PP_ALIGN.CENTER

        # Style header row differently
        for header_cell in table.rows[0].cells:
            header_cell.fill.solid()
            header_cell.fill.fore_color.rgb = style_config.background_color



class SlideManager:
    """
    A class to manage and arrange slides within a PowerPoint presentation.
    """

    def __init__(self, presentation):
        """
        Initialize the SlideManager with the given PowerPoint presentation.

        Args:
            presentation (Presentation): An existing PowerPoint presentation to manage.

        Raises:
            ValueError: If the presentation is not a valid Presentation object.
        """
        if not isinstance(presentation, Presentation):
            raise ValueError("Invalid presentation. Must be a PowerPoint Presentation object.")
        
        self.presentation = presentation

    def arrange_slides(self, order):
        """
        Arrange slides in the presentation according to the specified order.

        Args:
            order (list of int): List of indices representing the desired order of slides.

        Raises:
            ValueError: If the order list does not match the number of slides or indices are out of range.
        """
        if sorted(order) != list(range(len(self.presentation.slides))):
            raise ValueError("Order must be a permutation of all slide indices.")
        
        slides = list(self.presentation.slides)
        self.presentation.slides._spTree.clear()

        for index in order:
            self.presentation.slides._spTree.insert_element_before(
                slides[index].element, 'p:extLst'
            )


class StyleConfig:
    """
    A class to configure and validate styles for tables or text in PowerPoint slides.
    """

    def __init__(self, font_style, font_size, font_color, background_color):
        """
        Initialize the StyleConfig with specified styling attributes.

        Args:
            font_style (str): The font style to apply (e.g., 'Arial').
            font_size (int or float): The font size in points.
            font_color (tuple): RGB color tuple for the font (e.g., (0, 0, 0) for black).
            background_color (tuple): RGB color tuple for the background (e.g., (255, 255, 255) for white).

        Raises:
            ValueError: If any of the styling attributes are invalid.
        """
        self.font_style = font_style
        self.font_size = font_size
        self.font_color = font_color
        self.background_color = background_color

        self.validate_styles()

    def validate_styles(self):
        """
        Validate the provided styling attributes.

        Raises:
            ValueError: If font size is non-positive, colors are not valid RGB tuples, or the font style is not a string.
        """
        if not isinstance(self.font_style, str):
            raise ValueError("Font style must be a string.")
        
        if not isinstance(self.font_size, (int, float)) or self.font_size <= 0:
            raise ValueError("Font size must be a positive number.")
        
        if not isinstance(self.font_color, tuple) or len(self.font_color) != 3 or not all(0 <= c <= 255 for c in self.font_color):
            raise ValueError("Font color must be a tuple of three RGB values (0-255).")

        if not isinstance(self.background_color, tuple) or len(self.background_color) != 3 or not all(0 <= c <= 255 for c in self.background_color):
            raise ValueError("Background color must be a tuple of three RGB values (0-255).")




def load_data(file_path, **kwargs) -> pd.DataFrame:
    """
    Load data from a specified file path into a pandas DataFrame.

    Args:
        file_path (str): The path to the file containing the data.
        **kwargs: Optional keyword arguments for specific read functions.

    Returns:
        DataFrame: A pandas DataFrame containing the loaded data.

    Raises:
        ValueError: If the file format is unsupported or the file cannot be read.
        FileNotFoundError: If the file path does not exist.
    """
    if not isinstance(file_path, str) or not file_path:
        raise ValueError("File path must be a non-empty string.")
    
    try:
        if file_path.endswith('.csv'):
            return pd.read_csv(file_path, **kwargs)
        elif file_path.endswith(('.xls', '.xlsx')):
            return pd.read_excel(file_path, **kwargs)
        elif file_path.endswith('.json'):
            return pd.read_json(file_path, **kwargs)
        else:
            raise ValueError("Unsupported file format.")
    except FileNotFoundError:
        raise FileNotFoundError(f"The file at path {file_path} does not exist.")
    except Exception as e:
        raise ValueError(f"An error occurred while reading the file: {e}")




def style_table(table, **style_args):
    """
    Apply styling to a PowerPoint table.

    Args:
        table: The table object within a PowerPoint slide to which styles are applied.
        **style_args: Keyword arguments for styling attributes such as font style, font size, font color, etc.

    Raises:
        ValueError: If any style argument is invalid.
    """

    # Default styles
    default_styles = {
        "font_style": "Calibri",
        "font_size": 12,
        "font_color": (0, 0, 0),
        "background_color": (255, 255, 255),
        "border_style": 'thin', # Note: PowerPoint API limitations for border styling
        "alignment": 'center'
    }
    
    # Merge user-defined styles with defaults
    styles = {**default_styles, **style_args}

    # Validate RGB color tuples
    def validate_rgb(color):
        if not isinstance(color, tuple) or len(color) != 3 or not all(0 <= c <= 255 for c in color):
            raise ValueError("RGB color must be a tuple of three integers (0-255).")

    validate_rgb(styles['font_color'])
    validate_rgb(styles['background_color'])

    # Map string alignment to PP_ALIGN constants
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }

    if styles['alignment'] not in alignment_map:
        raise ValueError("Alignment must be 'left', 'center', or 'right'.")

    alignment = alignment_map[styles['alignment']]

    # Apply styles to each cell in the table
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                run = paragraph.runs[0]

                # Apply font style
                run.font.name = styles['font_style']
                # Apply font size
                run.font.size = Pt(styles['font_size'])
                # Apply font color
                run.font.color.rgb = RGBColor(*styles['font_color'])
                # Apply alignment
                paragraph.alignment = alignment

                # Set background color
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*styles['background_color'])



def customize_slide(slide, title, text_content, **style_args):
    """
    Customize a PowerPoint slide by setting the title, adding text content, and applying styles.

    Args:
        slide: The PowerPoint slide object to customize.
        title (str): The title to set in the slide's title placeholder.
        text_content (str): The main text content for the slide.
        **style_args: Keyword arguments for styling attributes such as font style, font size, font color, etc.

    Raises:
        ValueError: If any style argument is invalid or a required placeholder is missing.
    """
    # Default styles
    default_styles = {
        "font_style": "Calibri",
        "font_size": 12,
        "font_color": (0, 0, 0),
        "alignment": 'center',
        "placeholder_bg_color": (255, 255, 255)
    }
    
    # Merge user-defined styles with defaults
    styles = {**default_styles, **style_args}

    # Validate RGB color tuples
    def validate_rgb(color):
        if not isinstance(color, tuple) or len(color) != 3 or not all(0 <= c <= 255 for c in color):
            raise ValueError("RGB color must be a tuple of three integers (0-255).")

    validate_rgb(styles['font_color'])
    validate_rgb(styles['placeholder_bg_color'])

    # Map string alignment to PP_ALIGN constants
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT
    }

    if styles['alignment'] not in alignment_map:
        raise ValueError("Alignment must be 'left', 'center', or 'right'.")

    alignment = alignment_map[styles['alignment']]

    # Set the title
    if slide.shapes.title:
        slide.shapes.title.text = title
    
    # Add text box for main content
    left = Pt(1)
    top = Pt(2)
    width = slide.slide_width - Pt(2)
    height = slide.slide_height - Pt(4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text_content

    # Apply styles to text box
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = styles['font_style']
            run.font.size = Pt(styles['font_size'])
            run.font.color.rgb = RGBColor(*styles['font_color'])
            paragraph.alignment = alignment
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(*styles['placeholder_bg_color'])



def setup_logging(level):
    """
    Set up the logging configuration for the application.

    Args:
        level (str or int): The logging level to set. Valid levels are 'DEBUG', 'INFO',
                            'WARNING', 'ERROR', 'CRITICAL', or corresponding int values.

    Raises:
        ValueError: If the provided logging level is invalid.
    """
    # Define valid logging levels
    valid_levels = {
        'DEBUG': logging.DEBUG,
        'INFO': logging.INFO,
        'WARNING': logging.WARNING,
        'ERROR': logging.ERROR,
        'CRITICAL': logging.CRITICAL
    }

    # If level is a string, convert it to corresponding logging level
    if isinstance(level, str):
        level = level.upper()
        if level not in valid_levels:
            raise ValueError(f"Invalid logging level: '{level}'. Choose from {list(valid_levels.keys())}.")
        logging_level = valid_levels[level]
    elif isinstance(level, int):
        if level not in valid_levels.values():
            raise ValueError(f"Invalid logging level: {level}. Choose from {list(valid_levels.values())}.")
        logging_level = level
    else:
        raise ValueError("Logging level must be a string or integer.")

    # Configure the logging
    logging.basicConfig(
        level=logging_level,
        format='%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )

