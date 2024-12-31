from io import StringIO
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from your_module import DataFrameToSlide  # Replace 'your_module' with the actual module name
from your_module import PresentationGenerator  # Replace 'your_module' with the actual module name
from your_module import SlideManager  # Replace 'your_module' with the actual module name
from your_module import StyleConfig  # Replace 'your_module' with the actual module name
from your_module import customize_slide  # Replace 'your_module' with the actual module name
from your_module import load_data  # Replace 'your_module' with the actual module name
from your_module import setup_logging  # Replace 'your_module' with the actual module name
from your_module import style_table  # Replace 'your_module' with the actual module name
import logging
import os
import pandas as pd
import pytest


def test_init_with_valid_title():
    pg = PresentationGenerator("Sample Presentation")
    assert pg.title == "Sample Presentation"
    assert isinstance(pg.presentation, Presentation)
    assert len(pg.slides) == 1  # Title slide is added

def test_init_with_invalid_title():
    with pytest.raises(ValueError):
        PresentationGenerator(123)  # Title must be a string

def test_add_valid_text_slide():
    pg = PresentationGenerator("Test Presentation")
    pg.add_slide("Text Slide", "text", text_content="This is a text slide.")
    assert len(pg.slides) == 2
    assert pg.slides[1].shapes.title.text == "Text Slide"

def test_add_invalid_content_type():
    pg = PresentationGenerator("Test Presentation")
    with pytest.raises(ValueError):
        pg.add_slide("Invalid Slide", "image")  # Unsupported content type

def test_save_presentation_invalid_path():
    pg = PresentationGenerator("Test Presentation")
    with pytest.raises(ValueError):
        pg.save_presentation("invalid_path.txt")  # Must be .pptx file

def test_save_presentation_successful(tmpdir):
    pg = PresentationGenerator("Test Presentation")
    temp_file = os.path.join(tmpdir, "presentation.pptx")
    pg.save_presentation(temp_file)
    assert os.path.exists(temp_file)



# Mock StyleConfig class for testing
class StyleConfig:
    def __init__(self, font_style, font_size, font_color, background_color):
        self.font_style = font_style
        self.font_size = font_size
        self.font_color = font_color
        self.background_color = background_color

@pytest.fixture
def sample_dataframe():
    return pd.DataFrame({
        "Column1": ["Row1", "Row2"],
        "Column2": ["Value1", "Value2"]
    })

@pytest.fixture
def sample_style_config():
    return StyleConfig(
        font_style='Arial',
        font_size=Pt(12),
        font_color=RGBColor(0, 0, 0),
        background_color=RGBColor(255, 255, 255)
    )

@pytest.fixture
def slide():
    pres = Presentation()
    layout = pres.slide_layouts[5]  # Use blank slide layout
    return pres.slides.add_slide(layout)

def test_init_with_valid_dataframe(sample_dataframe, sample_style_config):
    df_to_slide = DataFrameToSlide(sample_dataframe, sample_style_config)
    assert df_to_slide.dataframe.equals(sample_dataframe)
    assert df_to_slide.style_config == sample_style_config

def test_init_with_invalid_dataframe(sample_style_config):
    with pytest.raises(ValueError):
        DataFrameToSlide("Not a DataFrame", sample_style_config)

def test_create_table_slide_valid(slide, sample_dataframe, sample_style_config):
    df_to_slide = DataFrameToSlide(sample_dataframe, sample_style_config)
    df_to_slide.create_table_slide(slide)
    assert len(slide.shapes) > 0  # Check that a table was added

def test_create_table_slide_invalid(slide, sample_style_config):
    with pytest.raises(ValueError):
        presentation = Presentation()  # Pass invalid slide object
        pd.DataFrameToSlide(presentation, sample_style_config).create_table_slide("Not a slide")

def test_apply_styles(sample_dataframe, sample_style_config, slide):
    df_to_slide = DataFrameToSlide(sample_dataframe, sample_style_config)
    df_to_slide.create_table_slide(slide)
    
    # Retrieve table from slide
    table = slide.shapes[0].table
    
    # Assertions to check if styles are applied
    assert table.cell(0, 0).text_frame.paragraphs[0].font.name == sample_style_config.font_style
    assert table.cell(0, 0).text_frame.paragraphs[0].font.size == sample_style_config.font_size
    assert table.cell(0, 0).text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT



@pytest.fixture
def sample_presentation():
    pres = Presentation()
    # Add three slides for testing
    pres.slides.add_slide(pres.slide_layouts[5])
    pres.slides.add_slide(pres.slide_layouts[5])
    pres.slides.add_slide(pres.slide_layouts[5])
    return pres

def test_init_with_valid_presentation(sample_presentation):
    manager = SlideManager(sample_presentation)
    assert manager.presentation == sample_presentation

def test_init_with_invalid_presentation():
    with pytest.raises(ValueError):
        SlideManager("Not a Presentation Object")

def test_arrange_slides_valid(sample_presentation):
    manager = SlideManager(sample_presentation)
    # Initial order should be [0, 1, 2]
    manager.arrange_slides([2, 0, 1])
    assert [slide.shapes.title.text for slide in manager.presentation.slides] == ['Slide 3', 'Slide 1', 'Slide 2']

def test_arrange_slides_invalid_order(sample_presentation):
    manager = SlideManager(sample_presentation)
    # Invalid order: doesn't match the number of slides
    with pytest.raises(ValueError):
        manager.arrange_slides([0, 2])

    # Invalid order: indices out of range
    with pytest.raises(ValueError):
        manager.arrange_slides([0, 1, 4])



def test_style_config_valid():
    config = StyleConfig('Arial', 12, (0, 0, 0), (255, 255, 255))
    assert config.font_style == 'Arial'
    assert config.font_size == 12
    assert config.font_color == (0, 0, 0)
    assert config.background_color == (255, 255, 255)

def test_style_config_invalid_font_style():
    with pytest.raises(ValueError, match="Font style must be a string."):
        StyleConfig(123, 12, (0, 0, 0), (255, 255, 255))

def test_style_config_invalid_font_size():
    with pytest.raises(ValueError, match="Font size must be a positive number."):
        StyleConfig('Arial', 0, (0, 0, 0), (255, 255, 255))

    with pytest.raises(ValueError, match="Font size must be a positive number."):
        StyleConfig('Arial', -5, (0, 0, 0), (255, 255, 255))

def test_style_config_invalid_font_color():
    with pytest.raises(ValueError, match="Font color must be a tuple of three RGB values \\(0-255\\)."):
        StyleConfig('Arial', 12, (256, 0, 0), (255, 255, 255))

    with pytest.raises(ValueError, match="Font color must be a tuple of three RGB values \\(0-255\\)."):
        StyleConfig('Arial', 12, (255, 0), (255, 255, 255))

def test_style_config_invalid_background_color():
    with pytest.raises(ValueError, match="Background color must be a tuple of three RGB values \\(0-255\\)."):
        StyleConfig('Arial', 12, (0, 0, 0), (255, -10, 255))

    with pytest.raises(ValueError, match="Background color must be a tuple of three RGB values \\(0-255\\)."):
        StyleConfig('Arial', 12, (0, 0, 0), (255,))



def test_load_data_csv(mocker):
    mock_csv_data = "col1,col2\nval1,val2\nval3,val4"
    mocker.patch("builtins.open", mocker.mock_open(read_data=mock_csv_data))
    mocker.patch("pandas.read_csv", return_value=pd.read_csv(StringIO(mock_csv_data)))
    df = load_data("test.csv")
    assert df.shape == (2, 2)
    assert list(df.columns) == ["col1", "col2"]

def test_load_data_excel(mocker):
    mock_excel_data = pd.DataFrame({'col1': ['val1', 'val2'], 'col2': ['val3', 'val4']})
    mocker.patch("pandas.read_excel", return_value=mock_excel_data)
    df = load_data("test.xlsx")
    assert df.shape == (2, 2)
    assert list(df.columns) == ["col1", "col2"]

def test_load_data_json(mocker):
    mock_json_data = '{"col1": ["val1", "val2"], "col2": ["val3", "val4"]}'
    mocker.patch("builtins.open", mocker.mock_open(read_data=mock_json_data))
    mocker.patch("pandas.read_json", return_value=pd.read_json(StringIO(mock_json_data)))
    df = load_data("test.json")
    assert df.shape == (2, 2)
    assert list(df.columns) == ["col1", "col2"]

def test_load_data_invalid_format():
    with pytest.raises(ValueError, match="Unsupported file format."):
        load_data("test.txt")

def test_load_data_file_not_found(mocker):
    mocker.patch("os.path.exists", return_value=False)
    with pytest.raises(FileNotFoundError):
        load_data("non_existing.csv")

def test_load_data_invalid_file_path():
    with pytest.raises(ValueError, match="File path must be a non-empty string."):
        load_data("")



@pytest.fixture
def sample_table():
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
    slide = presentation.slides.add_slide(slide_layout)
    table_shape = slide.shapes.add_table(2, 2, Pt(2), Pt(2), Pt(5), Pt(5))
    return table_shape.table

def test_style_table_default_styles(sample_table):
    style_table(sample_table)

    for row in sample_table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                run = paragraph.runs[0]
                assert run.font.name == 'Calibri'
                assert run.font.size == Pt(12)
                assert run.font.color.rgb == RGBColor(0, 0, 0)
                assert paragraph.alignment == PP_ALIGN.CENTER
                assert cell.fill.fore_color.rgb == RGBColor(255, 255, 255)

def test_style_table_custom_styles(sample_table):
    style_table(sample_table, font_style='Arial', font_size=14, font_color=(255, 0, 0), background_color=(0, 255, 0), alignment='right')

    for row in sample_table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                run = paragraph.runs[0]
                assert run.font.name == 'Arial'
                assert run.font.size == Pt(14)
                assert run.font.color.rgb == RGBColor(255, 0, 0)
                assert paragraph.alignment == PP_ALIGN.RIGHT
                assert cell.fill.fore_color.rgb == RGBColor(0, 255, 0)

def test_style_table_invalid_rgb():
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    table_shape = slide.shapes.add_table(2, 2, Pt(2), Pt(2), Pt(5), Pt(5))
    table = table_shape.table

    with pytest.raises(ValueError, match="RGB color must be a tuple of three integers \(0-255\)."):
        style_table(table, font_color=(256, 0, 0))

    with pytest.raises(ValueError, match="RGB color must be a tuple of three integers \(0-255\)."):
        style_table(table, background_color=(255, 255))

def test_style_table_invalid_alignment():
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    table_shape = slide.shapes.add_table(2, 2, Pt(2), Pt(2), Pt(5), Pt(5))
    table = table_shape.table

    with pytest.raises(ValueError, match="Alignment must be 'left', 'center', or 'right'."):
        style_table(table, alignment='bottom')



@pytest.fixture
def sample_slide():
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
    return presentation.slides.add_slide(slide_layout)

def test_customize_slide_default_styles(sample_slide):
    customize_slide(sample_slide, title="Sample Title", text_content="Sample content")
    
    title_shape = sample_slide.shapes.title
    if title_shape:
        assert title_shape.text == "Sample Title"

    text_box = next((shape for shape in sample_slide.shapes if shape.has_text_frame), None)
    text_frame = text_box.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            assert run.font.name == 'Calibri'
            assert run.font.size == Pt(12)
            assert run.font.color.rgb == RGBColor(0, 0, 0)
            assert paragraph.alignment == PP_ALIGN.CENTER
            assert text_box.fill.fore_color.rgb == RGBColor(255, 255, 255)

def test_customize_slide_custom_styles(sample_slide):
    customize_slide(sample_slide, title="Custom Title", text_content="Custom content", font_style="Arial", font_size=14, font_color=(255, 0, 0), alignment='right', placeholder_bg_color=(0, 255, 0))
    
    text_box = next((shape for shape in sample_slide.shapes if shape.has_text_frame), None)
    text_frame = text_box.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            assert run.font.name == 'Arial'
            assert run.font.size == Pt(14)
            assert run.font.color.rgb == RGBColor(255, 0, 0)
            assert paragraph.alignment == PP_ALIGN.RIGHT
            assert text_box.fill.fore_color.rgb == RGBColor(0, 255, 0)

def test_customize_slide_invalid_rgb(sample_slide):
    with pytest.raises(ValueError, match="RGB color must be a tuple of three integers \\(0-255\\)."):
        customize_slide(sample_slide, title="Title", text_content="Content", font_color=(256, 0, 0))

    with pytest.raises(ValueError, match="RGB color must be a tuple of three integers \\(0-255\\)."):
        customize_slide(sample_slide, title="Title", text_content="Content", placeholder_bg_color=(255, 0))

def test_customize_slide_invalid_alignment(sample_slide):
    with pytest.raises(ValueError, match="Alignment must be 'left', 'center', or 'right'."):
        customize_slide(sample_slide, title="Title", text_content="Content", alignment='bottom')



def test_setup_logging_valid_string_levels(caplog):
    # Test valid string levels
    for level in ['DEBUG', 'INFO', 'WARNING', 'ERROR', 'CRITICAL']:
        setup_logging(level)
        assert logging.getLogger().getEffectiveLevel() == logging._nameToLevel[level]

def test_setup_logging_valid_int_levels(caplog):
    # Test valid int levels
    for level in [logging.DEBUG, logging.INFO, logging.WARNING, logging.ERROR, logging.CRITICAL]:
        setup_logging(level)
        assert logging.getLogger().getEffectiveLevel() == level

def test_setup_logging_invalid_string_level():
    with pytest.raises(ValueError, match="Invalid logging level: 'INVALID'"):
        setup_logging("INVALID")

def test_setup_logging_invalid_int_level():
    with pytest.raises(ValueError, match="Invalid logging level: 100"):
        setup_logging(100)

def test_setup_logging_invalid_type():
    with pytest.raises(ValueError, match="Logging level must be a string or integer."):
        setup_logging(3.5)

@pytest.fixture
def reset_logging():
    """
    Fixture to reset the logging configuration before and after tests.
    This is necessary because logging.basicConfig can only be configured once in the current process runtime.
    """
    logging.shutdown()
    import importlib
    importlib.reload(logging)

def test_logging_output(caplog, reset_logging):
    setup_logging('DEBUG')
    with caplog.at_level(logging.DEBUG):
        logging.debug('Test debug message')
        logging.info('Test info message')

    assert 'DEBUG' in caplog.text
    assert 'Test debug message' in caplog.text
    assert 'INFO' in caplog.text
    assert 'Test info message' in caplog.text
