from DataFrameToPPT import DataFrameToPPT  # Assuming the class is in DataFrameToPPT.py
from PPTCustomizer import PPTCustomizer  # Assuming the class is in PPTCustomizer.py
from PPTManager import PPTManager  # Assuming the class is in PPTManager.py
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from setup_logging import setup_logging  # Assuming the function is in setup_logging.py
from validate_dataframe import validate_dataframe  # Assuming the function is in validate_dataframe.py
import logging
import os
import pandas as pd
import pytest


@pytest.fixture
def sample_dataframe():
    data = {
        'A': [1, 2, 3],
        'B': [4, 5, 6],
        'C': [7, 8, 9]
    }
    return pd.DataFrame(data, index=['X', 'Y', 'Z'])

@pytest.fixture
def sample_presentation():
    return Presentation()

def test_init(sample_dataframe):
    converter = DataFrameToPPT(sample_dataframe)
    assert converter.dataframe.equals(sample_dataframe)

def test_convert_to_bar_chart_creates_chart(sample_dataframe, sample_presentation):
    slide = sample_presentation.slides.add_slide(sample_presentation.slide_layouts[5])  # Title slide layout
    converter = DataFrameToPPT(sample_dataframe)
    chart = converter.convert_to_bar_chart(slide)
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED

def test_convert_to_line_chart_creates_chart(sample_dataframe, sample_presentation):
    slide = sample_presentation.slides.add_slide(sample_presentation.slide_layouts[5])
    converter = DataFrameToPPT(sample_dataframe)
    chart = converter.convert_to_line_chart(slide)
    assert chart.chart_type == XL_CHART_TYPE.LINE

def test_convert_to_pie_chart_creates_chart(sample_dataframe, sample_presentation):
    slide = sample_presentation.slides.add_slide(sample_presentation.slide_layouts[5])
    converter = DataFrameToPPT(sample_dataframe)
    chart = converter.convert_to_pie_chart(slide)
    assert chart.chart_type == XL_CHART_TYPE.PIE

def test_update_chart_title(sample_dataframe, sample_presentation):
    slide = sample_presentation.slides.add_slide(sample_presentation.slide_layouts[5])
    converter = DataFrameToPPT(sample_dataframe)
    chart = converter.convert_to_bar_chart(slide, title="Sample Title")
    assert chart.has_title
    assert chart.chart_title.text_frame.text == "Sample Title"

def test_update_chart_colors(sample_dataframe, sample_presentation):
    from pptx.dml.color import RGBColor
    slide = sample_presentation.slides.add_slide(sample_presentation.slide_layouts[5])
    converter = DataFrameToPPT(sample_dataframe)
    color_scheme = [RGBColor(0xFF, 0x00, 0x00), RGBColor(0x00, 0xFF, 0x00), RGBColor(0x00, 0x00, 0xFF)]
    chart = converter.convert_to_bar_chart(slide, color_scheme=color_scheme)
    for idx, series in enumerate(chart.series):
        assert len(color_scheme) == 3
        fill = series.format.fill
        fill_color = fill.fore_color.rgb
        assert fill_color == color_scheme[idx]



@pytest.fixture
def sample_presentation():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = ['A', 'B', 'C']
    chart_data.add_series('Series 1', (1, 2, 3))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(2), Inches(2), Inches(5), Inches(4), chart_data
    ).chart
    return chart

def test_set_title(sample_presentation):
    customizer = PPTCustomizer(sample_presentation)
    customizer.set_title("Sample Title", font_size=Pt(14), bold=True)
    assert sample_presentation.has_title
    assert sample_presentation.chart_title.text_frame.text == "Sample Title"
    assert sample_presentation.chart_title.text_frame.paragraphs[0].font.size == Pt(14)
    assert sample_presentation.chart_title.text_frame.paragraphs[0].font.bold is True

def test_set_axes_labels(sample_presentation):
    customizer = PPTCustomizer(sample_presentation)
    customizer.set_axes_labels("X Axis Label", "Y Axis Label", font_size=Pt(12))
    category_axis = sample_presentation.category_axis
    value_axis = sample_presentation.value_axis
    assert category_axis.axis_title.text_frame.text == "X Axis Label"
    assert value_axis.axis_title.text_frame.text == "Y Axis Label"
    assert category_axis.axis_title.text_frame.paragraphs[0].font.size == Pt(12)
    assert value_axis.axis_title.text_frame.paragraphs[0].font.size == Pt(12)

def test_set_legend(sample_presentation):
    customizer = PPTCustomizer(sample_presentation)
    customizer.set_legend(True, position=1)
    assert sample_presentation.has_legend is True
    assert sample_presentation.legend.position == 1

def test_set_colors(sample_presentation):
    customizer = PPTCustomizer(sample_presentation)
    color_scheme = [(255, 0, 0), (0, 255, 0), (0, 0, 255)]
    customizer.set_colors(color_scheme)
    for idx, series in enumerate(sample_presentation.series):
        fill_color = series.format.fill.fore_color.rgb
        assert fill_color == RGBColor(*color_scheme[idx])



@pytest.fixture
def ppt_manager(tmp_path):
    # Create a temporary presentation path
    temp_presentation_path = tmp_path / "test_presentation.pptx"
    return PPTManager(str(temp_presentation_path))

def test_init_creates_new_presentation(tmp_path):
    temp_presentation_path = tmp_path / "non_existent_presentation.pptx"
    manager = PPTManager(str(temp_presentation_path))
    assert isinstance(manager.presentation, Presentation)
    assert len(manager.presentation.slides) == 0

def test_add_slide(ppt_manager):
    slide_count_before = len(ppt_manager.presentation.slides)
    slide = ppt_manager.add_slide(0)  # Using the first layout as an example
    slide_count_after = len(ppt_manager.presentation.slides)
    assert slide_count_after == slide_count_before + 1
    assert slide in ppt_manager.presentation.slides

def test_save_presentation(ppt_manager, tmp_path):
    # Save presentation in a new path
    new_path = tmp_path / "new_test_presentation.pptx"
    ppt_manager.save_presentation(str(new_path))
    assert os.path.exists(new_path)

    # Modify and save in original path
    ppt_manager.add_slide(0)
    ppt_manager.save_presentation()
    original_path_exists = os.path.exists(ppt_manager.presentation_path)
    assert original_path_exists
    # Ensure that the modification (adding a slide) is saved
    assert len(Presentation(ppt_manager.presentation_path).slides) > 0



def test_valid_dataframe():
    df = pd.DataFrame({'A': [1, 2], 'B': [3, 4]})
    assert validate_dataframe(df) is True

def test_invalid_input_type():
    with pytest.raises(TypeError):
        validate_dataframe([1, 2, 3])  # Not a DataFrame

def test_empty_dataframe():
    df = pd.DataFrame()
    with pytest.raises(ValueError, match="DataFrame is empty and cannot be used for charting."):
        validate_dataframe(df)

def test_non_numeric_data():
    df = pd.DataFrame({'A': [1, 2], 'B': ['x', 'y']})
    with pytest.raises(ValueError, match="DataFrame contains non-numeric data."):
        validate_dataframe(df)

def test_dataframe_with_nan_values():
    df = pd.DataFrame({'A': [1, None], 'B': [3, 4]})
    with pytest.raises(ValueError, match="DataFrame contains NaN or missing values."):
        validate_dataframe(df)

def test_dataframe_with_zero_dimensions():
    df = pd.DataFrame([], columns=['A', 'B'])
    with pytest.raises(ValueError, match="DataFrame must have at least one row and one column."):
        validate_dataframe(df)



def test_setup_logging_invalid_level():
    with pytest.raises(ValueError, match="Invalid logging level. Must be one of DEBUG, INFO, WARNING, ERROR, CRITICAL."):
        setup_logging(level=99)  # Invalid logging level

def test_setup_logging_console(caplog):
    setup_logging(level=logging.DEBUG)
    
    with caplog.at_level(logging.DEBUG):
        logging.debug("This is a debug message.")

    assert "This is a debug message." in caplog.text

def test_setup_logging_file(tmp_path):
    log_file_path = tmp_path / "test_log.log"
    setup_logging(level=logging.INFO, log_file=str(log_file_path))

    logging.info("This is a test log entry.")
    logging.warning("This is a warning message.")

    with open(log_file_path, "r") as log_file:
        log_contents = log_file.read()

    assert "This is a test log entry." in log_contents
    assert "This is a warning message." in log_contents
