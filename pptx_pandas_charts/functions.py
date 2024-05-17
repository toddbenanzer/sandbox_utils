
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def dataframe_to_bar_chart(dataframe, title):
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    chart_data = CategoryChartData()
    for column in dataframe.columns:
        chart_data.categories.append(column)
    
    for index, row in dataframe.iterrows():
        chart_data.add_series(index, row.tolist())
    
    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5)
    slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data).chart
    
    return presentation

def dataframe_to_line_chart(dataframe, chart_title):
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    x_values = dataframe.index.tolist()
    y_values_list = [dataframe[col].tolist() for col in dataframe.columns]

    chart_data = CategoryChartData()
    chart_data.categories = x_values

    for idx, col in enumerate(dataframe.columns):
        chart_data.add_series(col, y_values_list[idx])

    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title

    return prs

def dataframe_to_scatter_plot(dataframe, chart_title):
    import matplotlib.pyplot as plt

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    fig, ax = plt.subplots()
    
    if 'x' in dataframe.columns and 'y' in dataframe.columns:
        dataframe.plot.scatter(x='x', y='y', ax=ax)
        ax.set_xlabel('x')
        ax.set_ylabel('y')
        ax.set_title(chart_title)
        
        temp_file = 'scatter_plot.png'
        plt.savefig(temp_file)
        
        slide.shapes.add_picture(temp_file, Inches(1), Inches(1), width=Inches(6), height=Inches(4))
        
        plt.close()

        return prs
    else:
        raise ValueError("DataFrame must contain 'x' and 'y' columns")

def convert_dataframe_to_pie_chart(df, chart_title):
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    x_categories = df.columns.tolist()
    
    values_list = [df[col].tolist() for col in df.columns]
    
   # Create the chart data object and populate it with data
   # Create the CategoryChartData object and populate it with data
    
   # Create the CategoryChartData object and populate it with data
   # Add the categories from the DataFrame columns to the CategoryChartData object
    
   # Add the values from the DataFrame rows to the CategoryChartData object
    
   # Create a new pie chart on the slide using the CategoryChartData object
    
   # Set the title of the pie chart
    
   # Return the PowerPoint presentation object
    
# Function to convert a DataFrame to an area chart

# Function to convert a DataFrame to a stacked bar chart

# Function to convert a DataFrame to a stacked area chart

# Function to convert a DataFrame to a stacked line chart


# Example usage of some functions:

if __name__ == "__main__":
  data_bar_chart_example()
  data_line_chart_example()
  data_scatter_plot_example()
  data_pie_chart_example()

def customize_x_axis_label(chart,x_axis_label,font_size=10,font_color=(0x00 0x00 0X00)):
      """
   Customize x-axis labels of a given Powerpoint Chart Object
   
      """
  
       # Accessing category axis of powerpoint Chart Object
    
       # Setting Title Text
       
       Setting Font Size
       
       Setting Font Color
        
       Setting Title Text
        
 
def customize_y_axis_label(chart,y_axis_label,font_size=12):
"""
Customize y-axis labels of powerpoint Chart Object.
"""
      # Accessing category axis of powerpoint Chart Object
      
       Setting Title Text
      
      Setting Font Size
 
