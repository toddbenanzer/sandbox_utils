from DataFrameToPPT import DataFrameToPPT  # Assuming the class is in DataFrameToPPT.py
from PPTCustomizer import PPTCustomizer  # Assuming the class is in PPTCustomizer.py
from PPTManager import PPTManager  # Assuming the class is in PPTManager.py
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Inches, Pt
from setup_logging import setup_logging  # Assuming the function is in setup_logging.py
from validate_dataframe import validate_dataframe  # Assuming the function is in validate_dataframe.py
import logging
import pandas as pd


# Example DataFrame
data = {
    'Product A': [23, 45, 56],
    'Product B': [19, 40, 25]
}
df = pd.DataFrame(data, index=['January', 'February', 'March'])

# Create a PowerPoint presentation
presentation = Presentation()
slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Title slide layout

# Example 1: Create a Bar Chart
df_to_ppt = DataFrameToPPT(df)
bar_chart = df_to_ppt.convert_to_bar_chart(slide, title="Monthly Sales Data", color_scheme=[
    RGBColor(0x00, 0x00, 0xFF), 
    RGBColor(0xFF, 0x00, 0x00)
])

# Example 2: Create a Line Chart
line_chart = df_to_ppt.convert_to_line_chart(slide, title="Monthly Sales Data (Line)", color_scheme=[
    RGBColor(0x00, 0xFF, 0x00), 
    RGBColor(0x00, 0x00, 0xFF)
])

# Example 3: Create a Pie Chart
slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])  # Add another slide
df_single_series = df.sum(axis=1).to_frame('Total Sales')  # Aggregate data for a pie chart
df_to_ppt_single = DataFrameToPPT(df_single_series)
pie_chart = df_to_ppt_single.convert_to_pie_chart(slide2, title="Sales Distribution")

# Save the presentation
presentation_path = 'monthly_sales_report.pptx'
presentation.save(presentation_path)

print(f"Presentation saved as {presentation_path}")



# Create a sample presentation and add a slide
presentation = Presentation()
slide = presentation.slides.add_slide(presentation.slide_layouts[5])

# Define chart data
chart_data = CategoryChartData()
chart_data.categories = ['A', 'B', 'C']
chart_data.add_series('Series 1', (5, 7, 3))
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(2), Inches(2), Inches(5), Inches(4), chart_data
)
chart = chart_shape.chart

# Example 1: Set Chart Title
customizer = PPTCustomizer(chart)
customizer.set_title("Sales Data", font_size=Pt(18), bold=True)

# Example 2: Set Axes Labels
customizer.set_axes_labels("Product Category", "Sales", font_size=Pt(12))

# Example 3: Set Legend Display
customizer.set_legend(True, position=1)

# Example 4: Set Colors for Series
color_scheme = [(255, 0, 0), (0, 255, 0), (0, 0, 255)]
customizer.set_colors(color_scheme)

# Save the presentation
presentation.save('customized_chart_presentation.pptx')



# Example 1: Creating a new presentation and adding a slide
presentation_path = "example_presentation.pptx"
manager = PPTManager(presentation_path)
slide = manager.add_slide(0)  # Add a new slide using the first layout
manager.save_presentation()

# Example 2: Adding a second slide to an existing presentation
slide2 = manager.add_slide(1)  # Add another slide using a different layout
manager.save_presentation()

# Example 3: Saving to a different file path
alternative_path = "alternative_presentation.pptx"
manager.save_presentation(alternative_path)



# Example 1: Valid DataFrame with numeric data
try:
    df_valid = pd.DataFrame({'Column1': [1, 2, 3], 'Column2': [4, 5, 6]})
    is_valid = validate_dataframe(df_valid)
    print("Valid DataFrame:", is_valid)  # Output: Valid DataFrame: True
except Exception as e:
    print(e)

# Example 2: DataFrame with non-numeric data
try:
    df_non_numeric = pd.DataFrame({'Column1': [1, 2, 3], 'Column2': ['a', 'b', 'c']})
    validate_dataframe(df_non_numeric)
except Exception as e:
    print("Error:", e)  # Output: Error: DataFrame contains non-numeric data.

# Example 3: DataFrame with NaN values
try:
    df_with_nan = pd.DataFrame({'Column1': [1, 2, None], 'Column2': [4, 5, 6]})
    validate_dataframe(df_with_nan)
except Exception as e:
    print("Error:", e)  # Output: Error: DataFrame contains NaN or missing values.

# Example 4: Empty DataFrame
try:
    df_empty = pd.DataFrame()
    validate_dataframe(df_empty)
except Exception as e:
    print("Error:", e)  # Output: Error: DataFrame is empty and cannot be used for charting.



# Example 1: Basic console logging setup at DEBUG level
setup_logging(level=logging.DEBUG)

# Log messages at various levels
logging.debug("Debugging information")
logging.info("General information")
logging.warning("Warning message")
logging.error("Error message")
logging.critical("Critical error message")

# Example 2: Logging to a file and console with INFO level
setup_logging(level=logging.INFO, log_file="app_log.txt")

# Log messages at various levels
logging.info("This is an informational message written to both console and file.")
logging.error("This error message will also appear in the log file.")

# Example 3: Invalid logging level demonstration
try:
    setup_logging(level=99)
except ValueError as e:
    print("Error setting up logging:", e)  # Output: Error setting up logging: Invalid logging level. Must be one of DEBUG, INFO, WARNING, ERROR, CRITICAL.
