from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.plot import Plot
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import logging
import pandas as pd


class DataFrameToPPT:
    """
    A class to convert pandas DataFrames to PowerPoint charts.
    """

    def __init__(self, dataframe: pd.DataFrame):
        """
        Initializes the DataFrameToPPT instance with the specified DataFrame.

        Args:
            dataframe (pd.DataFrame): The source data to be converted into charts.
        """
        self.dataframe = dataframe

    def convert_to_chart(self, slide, chart_type: XL_CHART_TYPE, **kwargs):
        """
        General function to convert DataFrame data into specified chart type on a PowerPoint slide.

        Args:
            slide (pptx.slide.Slide): The slide to which the chart will be added.
            chart_type (XL_CHART_TYPE): The type of chart to be created.
            **kwargs: Additional parameters for chart customization.

        Returns:
            pptx.chart.Chart: The created chart object.
        """
        chart_data = CategoryChartData()
        chart_data.categories = self.dataframe.index

        for column in self.dataframe.columns:
            chart_data.add_series(column, self.dataframe[column])

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data
        ).chart

        self.update_chart(chart, **kwargs)

        return chart

    def convert_to_bar_chart(self, slide, **kwargs):
        """
        Converts DataFrame data into a bar chart on the specified slide.

        Args:
            slide (pptx.slide.Slide): The slide to which the bar chart will be added.
            **kwargs: Additional parameters for chart customization.

        Returns:
            pptx.chart.Chart: The created bar chart object.
        """
        return self.convert_to_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, **kwargs)

    def convert_to_line_chart(self, slide, **kwargs):
        """
        Converts DataFrame data into a line chart on the specified slide.

        Args:
            slide (pptx.slide.Slide): The slide to which the line chart will be added.
            **kwargs: Additional parameters for chart customization.

        Returns:
            pptx.chart.Chart: The created line chart object.
        """
        return self.convert_to_chart(slide, XL_CHART_TYPE.LINE, **kwargs)

    def convert_to_pie_chart(self, slide, **kwargs):
        """
        Converts DataFrame data into a pie chart on the specified slide.

        Args:
            slide (pptx.slide.Slide): The slide to which the pie chart will be added.
            **kwargs: Additional parameters for chart customization.

        Returns:
            pptx.chart.Chart: The created pie chart object.
        """
        return self.convert_to_chart(slide, XL_CHART_TYPE.PIE, **kwargs)

    def update_chart(self, chart, **kwargs):
        """
        Updates an existing PowerPoint chart with new data or customization.

        Args:
            chart (pptx.chart.Chart): The PowerPoint chart to be updated.
            **kwargs: Various update parameters (new data, title, colors, etc.)

        Returns:
            None
        """
        if 'title' in kwargs:
            chart.has_title = True
            chart.chart_title.text_frame.text = kwargs['title']

        if 'color_scheme' in kwargs:
            # Assume color_scheme is a list of RGB tuples corresponding to series color
            for idx, series in enumerate(chart.series):
                if idx < len(kwargs['color_scheme']):
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = kwargs['color_scheme'][idx]



class PPTCustomizer:
    """
    A class to customize PowerPoint chart objects.
    """

    def __init__(self, chart):
        """
        Initializes the PPTCustomizer with the specified PowerPoint chart object.
        
        Args:
            chart (pptx.chart.Chart): The PowerPoint chart to be customized.
        """
        self.chart = chart

    def set_title(self, title, **kwargs):
        """
        Sets or updates the title of the specified chart.
        
        Args:
            title (str): The title text to be set on the chart.
            **kwargs: Additional customization parameters for the title (e.g., font size, boldness).
        """
        self.chart.has_title = True
        self.chart.chart_title.text_frame.text = title

        # Example of additional customization
        if 'font_size' in kwargs:
            self.chart.chart_title.text_frame.paragraphs[0].font.size = kwargs['font_size']
        
        if 'bold' in kwargs:
            self.chart.chart_title.text_frame.paragraphs[0].font.bold = kwargs['bold']

    def set_axes_labels(self, x_label, y_label, **kwargs):
        """
        Sets the labels for the x-axis and y-axis on the specified chart.
        
        Args:
            x_label (str): Label for the x-axis.
            y_label (str): Label for the y-axis.
            **kwargs: Additional customization parameters (e.g., font, size, rotation).
        """
        category_axis = self.chart.category_axis
        value_axis = self.chart.value_axis
        category_axis.has_title = True
        value_axis.has_title = True
        category_axis.axis_title.text_frame.text = x_label
        value_axis.axis_title.text_frame.text = y_label

        # Example of customization for axis labels
        if 'font_size' in kwargs:
            category_axis.axis_title.text_frame.paragraphs[0].font.size = kwargs['font_size']
            value_axis.axis_title.text_frame.paragraphs[0].font.size = kwargs['font_size']

    def set_legend(self, display, **kwargs):
        """
        Configures the legend display settings for the specified chart.
        
        Args:
            display (bool): Flag to show or hide the legend.
            **kwargs: Additional customization parameters for the legend (e.g., position, font).
        """
        self.chart.has_legend = display

        # Example of additional customization
        if display and 'position' in kwargs:
            self.chart.legend.position = kwargs['position']

    def set_colors(self, color_scheme):
        """
        Applies a color scheme to the series in the specified chart.
        
        Args:
            color_scheme (list of RGB tuples): e.g., [(255, 0, 0), (0, 255, 0)] specifying colors for the series.
        """
        for idx, series in enumerate(self.chart.series):
            if idx < len(color_scheme):
                fill = series.format.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*color_scheme[idx])



class PPTManager:
    """
    A class to manage PowerPoint presentations.
    """

    def __init__(self, presentation_path: str):
        """
        Initializes the PPTManager with the specified path for a new or existing PowerPoint presentation.
        
        Args:
            presentation_path (str): Path to the PowerPoint file to open or create.
        """
        try:
            self.presentation = Presentation(presentation_path)
        except Exception as e:
            self.presentation = Presentation()
        self.presentation_path = presentation_path

    def add_slide(self, layout_index: int):
        """
        Adds a new slide with the specified layout to the PowerPoint presentation.
        
        Args:
            layout_index (int): Index of the layout in the presentation's slide layout collection.
        
        Returns:
            pptx.slide.Slide: The newly created slide object.
        """
        slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(slide_layout)
        return slide

    def save_presentation(self, path: str = None):
        """
        Saves the current PowerPoint presentation to the specified file path.
        
        Args:
            path (str, optional): Path to save the PowerPoint presentation. 
                                  If None, saves to the initial presentation path.
        
        Returns:
            None
        """
        save_path = path if path else self.presentation_path
        self.presentation.save(save_path)



def validate_dataframe(dataframe: pd.DataFrame) -> bool:
    """
    Validates the input DataFrame to ensure it is suitable for chart conversion.

    Args:
        dataframe (pd.DataFrame): The DataFrame to be validated.

    Returns:
        bool: True if the DataFrame is valid, False otherwise.

    Raises:
        TypeError: If the input is not a pandas DataFrame.
        ValueError: If the DataFrame is empty, has non-numeric data, or contains NaN values.
    """
    # Check if input is a DataFrame
    if not isinstance(dataframe, pd.DataFrame):
        raise TypeError("Input must be a pandas DataFrame.")

    # Check if DataFrame is empty
    if dataframe.empty:
        raise ValueError("DataFrame is empty and cannot be used for charting.")

    # Check for non-zero dimensions
    if dataframe.shape[0] == 0 or dataframe.shape[1] == 0:
        raise ValueError("DataFrame must have at least one row and one column.")

    # Check for numeric data
    if not all(dataframe.applymap(lambda x: isinstance(x, (int, float))).all()):
        raise ValueError("DataFrame contains non-numeric data.")

    # Check for NaN or missing values
    if dataframe.isnull().values.any():
        raise ValueError("DataFrame contains NaN or missing values.")

    return True



def setup_logging(level: int, log_file: str = None, **kwargs) -> None:
    """
    Configures the logging utility for the application.

    Args:
        level (int): The logging level to set. It determines the severity of messages to capture.
        log_file (str, optional): The path to the log file. If None, logs are written only to the console.
        **kwargs: Additional arguments for further customization.

    Raises:
        ValueError: If an invalid logging level is provided.
    """
    # Validate and set the logging level
    if level not in [logging.DEBUG, logging.INFO, logging.WARNING, logging.ERROR, logging.CRITICAL]:
        raise ValueError("Invalid logging level. Must be one of DEBUG, INFO, WARNING, ERROR, CRITICAL.")

    # Format for log messages
    log_format = '%(asctime)s - %(levelname)s - %(message)s'
    
    # Configure logging handlers
    handlers = [logging.StreamHandler()]
    
    if log_file:
        handlers.append(logging.FileHandler(log_file))
    
    logging.basicConfig(level=level, format=log_format, handlers=handlers, **kwargs)
