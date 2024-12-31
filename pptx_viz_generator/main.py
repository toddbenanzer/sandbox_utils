from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pptx.util import Inches
import io
import json
import os
import pandas as pd
import plotly.graph_objects as go


class PresentationBuilder:
    """
    A class to build and manage a PowerPoint presentation using pptx-python.
    """

    def __init__(self):
        """
        Initializes a new PresentationBuilder object, setting up an empty presentation.
        """
        self.presentation = Presentation()

    def add_slide(self, layout):
        """
        Adds a new slide to the presentation using the specified SlideLayout.

        Args:
            layout: A SlideLayout object that defines the layout for the new slide.

        Returns:
            slide: The newly added slide object.
        """
        slide = self.presentation.slides.add_slide(layout.apply_layout(self.presentation.slides.layout))
        return slide
    
    def save(self, file_path):
        """
        Saves the current state of the presentation to a specified file path.

        Args:
            file_path (str): The path where the presentation file will be saved.
        """
        self.presentation.save(file_path)



class SlideLayout:
    """
    A class to define and apply slide layouts with specified design elements.
    """

    def __init__(self, layout_name, font, color_scheme):
        """
        Initializes a new SlideLayout object with the given layout name, font, and color scheme.

        Args:
            layout_name (str): The name of the layout for identification.
            font (str): The font style to be used in the slide layout.
            color_scheme (str): The color scheme to be applied to the slide.
        """
        self.layout_name = layout_name
        self.font = font
        self.color_scheme = color_scheme

    def apply_layout(self, slide):
        """
        Applies the defined layout, font, and color scheme to the specified slide.

        Args:
            slide (Slide): The slide object to which the layout will be applied.

        Returns:
            Slide: The modified slide object with the specified layout and style.
        """
        # Example code to apply font and color scheme
        # As pptx does not have direct methods to set font and color scheme, this will be a stub
        # Implement actual changes to the slide such as changing placeholder text properties
        # Set font (illustrative, you need more to actually apply fonts in pptx)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = self.font
            # Color scheme would generally require setting colors for each elements
            # An example could be manually setting background or text colors based on the scheme

        return slide



class PlotlyChartEmbedder:
    """
    A class to embed Plotly charts in PowerPoint slides.
    """

    def __init__(self, chart_data):
        """
        Initializes a PlotlyChartEmbedder with the given chart data.

        Args:
            chart_data: Data necessary to create a Plotly chart, typically a Plotly figure.
        """
        self.chart_data = chart_data

    def embed_chart(self, slide, position):
        """
        Embeds a Plotly chart into a specified slide at a given position.

        Args:
            slide (Slide): The slide object where the chart will be embedded.
            position (tuple): A tuple (x, y, width, height) defining the chart's position and dimensions.

        Returns:
            Slide: The modified slide with the embedded chart.
        """
        x, y, width, height = position

        # Save the Plotly chart as an image (e.g., PNG) in memory
        with io.BytesIO() as image_stream:
            self.chart_data.write_image(image_stream, format='png')
            image_stream.seek(0)

            # Add the image to the slide
            slide.shapes.add_picture(image_stream, Inches(x), Inches(y), Inches(width), Inches(height))

        return slide



class BrandingFormatter:
    """
    A class to apply branding guidelines to a PowerPoint presentation.
    """
    
    def __init__(self, branding_guidelines):
        """
        Initializes a BrandingFormatter with the given branding guidelines.

        Args:
            branding_guidelines (dict): Dictionary containing branding elements like fonts and colors.
        """
        self.branding_guidelines = branding_guidelines

    def apply_branding(self, presentation):
        """
        Applies the branding guidelines to each slide and element in the presentation.

        Args:
            presentation (Presentation): The PowerPoint presentation object to apply branding to.

        Returns:
            Presentation: The modified presentation with applied branding.
        """
        font = self.branding_guidelines.get('font', 'Arial')
        color_scheme = self.branding_guidelines.get('color_scheme', '000000')  # Default to black

        # Iterate over all slides and apply branding
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font
                            run.font.color.rgb = color_scheme  # These operations depend on `pptx` library capability

        return presentation



def load_branding_guidelines(guidelines_file):
    """
    Loads and parses branding guidelines from a specified file.

    Args:
        guidelines_file (str): The file path to the branding guidelines in JSON format.

    Returns:
        dict: A dictionary containing the parsed branding elements like fonts and color schemes.

    Raises:
        FileNotFoundError: If the guidelines file is not found.
        ValueError: If there is an error parsing the file contents.
    """
    if not os.path.exists(guidelines_file):
        raise FileNotFoundError(f"The file '{guidelines_file}' does not exist.")

    try:
        with open(guidelines_file, 'r') as file:
            guidelines = json.load(file)
    except json.JSONDecodeError as e:
        raise ValueError(f"Error parsing the branding guidelines file: {e}")

    return guidelines



def create_custom_layout(layout_details):
    """
    Creates a custom slide layout based on specified details.

    Args:
        layout_details (dict): Dictionary containing layout configurations such as positions and style attributes.

    Returns:
        SlideLayout: A SlideLayout object configured according to the provided layout details.
    """
    # Extract layout details from the dictionary
    title_position = layout_details.get('title_position', (1, 1, 5, 1))
    content_position = layout_details.get('content_position', (1, 2, 5, 4))
    background_color = layout_details.get('background_color', 'FFFFFF')  # Default white background
    font_style = layout_details.get('font_style', {'family': 'Arial', 'size': 12, 'color': '000000'})  # Default to Arial
    
    # Create a customized SlideLayout instance
    class CustomSlideLayout(SlideLayout):
        def __init__(self):
            super().__init__(layout_name="Custom Layout", 
                             font=font_style['family'], 
                             color_scheme=RGBColor.from_string(font_style['color']))
            
        def apply_layout(self, slide):
            """
            Applies the defined custom layout to a specified slide.

            Args:
                slide (Slide): The slide object to which the layout will be applied.

            Returns:
                Slide: The modified slide object with the specified custom layout.
            """
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor.from_string(background_color)
            
            # Set title position and format
            title = slide.shapes.title
            title.text_frame.paragraphs[0].font.size = font_style['size']
            title.left, title.top, title.width, title.height = (Inches(pos) for pos in title_position)
            
            # Set content position
            content = slide.placeholders[1]
            content.left, content.top, content.width, content.height = (Inches(pos) for pos in content_position)
            
            return slide

    return CustomSlideLayout()



def generate_interactive_chart(data, chart_type):
    """
    Generates an interactive chart using Plotly based on the given data and chart type.

    Args:
        data: Data structure such as a DataFrame, dictionary, or list containing the data to be visualized.
        chart_type (str): Specifies the type of chart to generate (e.g., 'bar', 'line', 'scatter', 'pie').

    Returns:
        go.Figure: A Plotly Figure object representing the interactive chart.

    Raises:
        ValueError: If the chart type is unsupported or data format is not compatible.
    """
    if not isinstance(data, (pd.DataFrame, dict, list)):
        raise ValueError("Data should be a pandas DataFrame, dictionary, or list.")
    
    if chart_type == 'bar':
        if isinstance(data, pd.DataFrame):
            fig = go.Figure(data=[go.Bar(x=data.iloc[:, 0], y=data.iloc[:, 1])])
        elif isinstance(data, dict):
            fig = go.Figure(data=[go.Bar(x=list(data.keys()), y=list(data.values()))])
        else:
            raise ValueError("Unsupported data format for 'bar' chart.")

    elif chart_type == 'line':
        if isinstance(data, pd.DataFrame):
            fig = go.Figure(data=[go.Scatter(x=data.iloc[:, 0], y=data.iloc[:, 1], mode='lines')])
        elif isinstance(data, dict):
            fig = go.Figure(data=[go.Scatter(x=list(data.keys()), y=list(data.values()), mode='lines')])
        else:
            raise ValueError("Unsupported data format for 'line' chart.")

    elif chart_type == 'scatter':
        if isinstance(data, pd.DataFrame):
            fig = go.Figure(data=[go.Scatter(x=data.iloc[:, 0], y=data.iloc[:, 1], mode='markers')])
        elif isinstance(data, dict):
            fig = go.Figure(data=[go.Scatter(x=list(data.keys()), y=list(data.values()), mode='markers')])
        else:
            raise ValueError("Unsupported data format for 'scatter' chart.")

    elif chart_type == 'pie':
        if isinstance(data, pd.DataFrame) and data.shape[1] == 2:
            fig = go.Figure(data=[go.Pie(labels=data.iloc[:, 0], values=data.iloc[:, 1])])
        elif isinstance(data, dict):
            fig = go.Figure(data=[go.Pie(labels=list(data.keys()), values=list(data.values()))])
        else:
            raise ValueError("Unsupported data format for 'pie' chart. Requires two columns for labels and values.")

    else:
        raise ValueError(f"Unsupported chart type: {chart_type}")

    # Basic layout for the chart
    fig.update_layout(title=f"Generated {chart_type.title()} Chart")
    
    return fig
