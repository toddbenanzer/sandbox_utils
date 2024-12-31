from brand_formatter import BrandingFormatter  # Hypothetical import
from branding_formatter import BrandingFormatter  # Assuming this class is in branding_formatter.py
from create_custom_layout import create_custom_layout  # Adjust the import path according to your project structure
from generate_interactive_chart import generate_interactive_chart  # Adjust import path as needed
from plotly_chart_embedder import PlotlyChartEmbedder  # Assuming this class is in plotly_chart_embedder.py
from pptx import Presentation
from pptx.util import Inches
from presentation_builder import PresentationBuilder
from slide_layout import SlideLayout  # Assuming SlideLayout is in slide_layout.py
from slide_layout import SlideLayout  # Hypothetical import for SlideLayout class
import pandas as pd
import plotly.graph_objects as go


# Example usage 1: Creating a blank presentation
builder = PresentationBuilder()

# Adding a slide with a custom layout
custom_layout = SlideLayout(layout_name="Title and Content", font="Arial", color_scheme="Blue")
slide = builder.add_slide(custom_layout)

# Saving the presentation to a file
builder.save("example_presentation_1.pptx")

# Example usage 2: Creating a presentation with multiple slides
builder = PresentationBuilder()

# Adding the first slide
slide1_layout = SlideLayout(layout_name="Title Slide", font="Times New Roman", color_scheme="Red")
slide1 = builder.add_slide(slide1_layout)

# Adding the second slide
slide2_layout = SlideLayout(layout_name="Content Only", font="Calibri", color_scheme="Green")
slide2 = builder.add_slide(slide2_layout)

# Saving the presentation with two slides
builder.save("example_presentation_2.pptx")

# Example usage 3: Using BrandingFormatter hypothetically

# Initialize PresentationBuilder
builder = PresentationBuilder()

# Adding a branded slide
branding = BrandingFormatter(branding_guidelines={"font": "Verdana", "color_scheme": "Corporate"})
branded_layout = SlideLayout(layout_name="Branded Slide", font="Verdana", color_scheme="Corporate")
slide3 = builder.add_slide(branded_layout)

# Applying branding
branding.apply_branding(builder.presentation)

# Save the presentation
builder.save("example_presentation_3.pptx")



# Example usage 1: Applying a basic layout to a new slide
presentation = Presentation()
slide_layout = SlideLayout("Title Slide", "Arial", "Blue")

# Add a new blank slide
slide_layout_type = presentation.slide_layouts[5]  # Use a blank layout
new_slide = presentation.slides.add_slide(slide_layout_type)

# Apply the custom layout
slide_layout.apply_layout(new_slide)

# Save the modified presentation
presentation.save("basic_layout_presentation.pptx")

# Example usage 2: Creating a presentation with multiple slides with different layouts
presentation = Presentation()

# Slide 1: Title Slide Layout
title_layout = SlideLayout("Title", "Times New Roman", "Red")
title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title_layout.apply_layout(title_slide)

# Slide 2: Content Slide Layout
content_layout = SlideLayout("Content", "Verdana", "Green")
content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
content_layout.apply_layout(content_slide)

# Save the presentation with customized slides
presentation.save("multi_layout_presentation.pptx")



# Example usage 1: Embedding a bar chart in a new slide
# Create a Plotly bar chart
fig = go.Figure(data=go.Bar(y=[2, 3, 1], x=['A', 'B', 'C']))

# Instantiate PlotlyChartEmbedder with the chart data
chart_embedder = PlotlyChartEmbedder(fig)

# Create a new presentation and slide
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Embed the chart at specific position (1 inch from left, 1 inch from top, 4x3 inches size)
chart_embedder.embed_chart(slide, (1, 1, 4, 3))

# Save the presentation
presentation.save("embedded_chart_presentation.pptx")

# Example usage 2: Embedding multiple charts in different slides
# Create another Plotly chart, such as a line chart
line_chart = go.Figure(data=go.Scatter(x=[1, 2, 3], y=[4, 1, 2]))

# Create a new slide for the line chart
line_slide = presentation.slides.add_slide(slide_layout)

# Embed the line chart
line_embedder = PlotlyChartEmbedder(line_chart)
line_embedder.embed_chart(line_slide, (1, 1, 5, 4))

# Save the updated presentation
presentation.save("multiple_charts_presentation.pptx")



# Example usage 1: Applying simple branding to a presentation
# Create a new presentation
presentation = Presentation()

# Add a slide with a title and content layout
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Welcome to Our Presentation"

# Set the branding guidelines
branding_guidelines = {
    "font": "Calibri",
    "color_scheme": "FF5733"
}

# Create a BrandingFormatter and apply branding
formatter = BrandingFormatter(branding_guidelines)
formatter.apply_branding(presentation)

# Save the branded presentation
presentation.save("branded_presentation.pptx")

# Example usage 2: Applying branding to multiple slides
# Add another slide to the presentation
slide_layout = presentation.slide_layouts[1]
slide2 = presentation.slides.add_slide(slide_layout)
content = slide2.shapes.placeholders[1]
content.text = "Here is some important information."

# Apply branding again
formatter.apply_branding(presentation)

# Save the updated presentation
presentation.save("multi_slide_branded_presentation.pptx")


# Example usage 1: Load branding guidelines from a valid JSON file
try:
    guidelines = load_branding_guidelines("branding_guidelines.json")
    print(guidelines)
except FileNotFoundError as e:
    print(e)
except ValueError as e:
    print(e)

# Example usage 2: Handle missing file error
try:
    guidelines = load_branding_guidelines("non_existent_file.json")
    print(guidelines)
except FileNotFoundError as e:
    print(f"Error: {e}")

# Example usage 3: Handle invalid JSON content
try:
    # Assuming 'invalid.json' contains malformed JSON
    guidelines = load_branding_guidelines("invalid.json")
    print(guidelines)
except ValueError as e:
    print(f"Parsing Error: {e}")



# Example usage 1: Applying a custom layout with specific details
# Create a dictionary with custom layout details
layout_details = {
    'title_position': (1, 0.5, 8, 1),  # x, y, width, height in inches
    'content_position': (1, 2, 8, 4),
    'background_color': 'EFEFEF',
    'font_style': {'family': 'Calibri', 'size': 24, 'color': '002366'}
}

# Create a custom layout using the provided details
custom_layout_class = create_custom_layout(layout_details)
custom_layout = custom_layout_class()

# Create a presentation and add a slide to apply the layout
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Custom Layout Title"
slide.placeholders[1].text = "This is the main content area."

# Apply the defined custom layout
custom_layout.apply_layout(slide)

# Save the presentation
presentation.save("custom_layout_presentation.pptx")

# Example usage 2: Applying a custom layout with default values
# Create a layout with default configurations
default_layout_class = create_custom_layout({})
default_layout = default_layout_class()

# Add a slide using the default layout settings
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)
slide.shapes.title.text = "Default Layout Title"
slide.placeholders[1].text = "Content for default layout."

# Apply the default custom layout
default_layout.apply_layout(slide)

# Save the presentation
presentation.save("default_layout_presentation.pptx")



# Example usage 1: Creating a bar chart from a DataFrame
data = pd.DataFrame({
    'Category': ['A', 'B', 'C'],
    'Values': [10, 15, 13]
})
bar_chart = generate_interactive_chart(data, 'bar')
bar_chart.show()

# Example usage 2: Creating a line chart from a dictionary
data_dict = {'Jan': 100, 'Feb': 150, 'Mar': 200}
line_chart = generate_interactive_chart(data_dict, 'line')
line_chart.show()

# Example usage 3: Creating a scatter plot from a DataFrame
scatter_data = pd.DataFrame({
    'X': [1, 2, 3, 4, 5],
    'Y': [10, 11, 12, 13, 14]
})
scatter_chart = generate_interactive_chart(scatter_data, 'scatter')
scatter_chart.show()

# Example usage 4: Creating a pie chart using a dictionary
pie_data = {'Apples': 5, 'Bananas': 3, 'Cherries': 7}
pie_chart = generate_interactive_chart(pie_data, 'pie')
pie_chart.show()
