from branding_formatter import BrandingFormatter  # Assuming the class is in branding_formatter.py
from custom_layout import create_custom_layout  # Assuming the function is in custom_layout.py
from generate_interactive_chart import generate_interactive_chart  # Adjust import path as needed
from load_branding_guidelines import load_branding_guidelines  # Assume this function is in load_branding_guidelines.py
from plotly.graph_objs import Figure
from plotly_chart_embedder import PlotlyChartEmbedder  # Assuming the class is in plotly_chart_embedder.py
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches
from presentation_builder import PresentationBuilder
from slide_layout import SlideLayout  # Assuming SlideLayout is implemented
from slide_layout import SlideLayout  # Assuming SlideLayout is in slide_layout.py
from unittest import mock
from unittest.mock import Mock
import json
import pandas as pd
import plotly.graph_objects as go
import pytest


@pytest.fixture
def mock_layout():
    # Create a mock layout with an apply_layout method returning a layout index
    mock = Mock()
    mock.apply_layout.return_value = 0
    return mock

@pytest.fixture
def presentation_builder():
    # Instantiate the PresentationBuilder object
    return PresentationBuilder()

def test_init_sets_empty_presentation(presentation_builder):
    assert isinstance(presentation_builder.presentation, Presentation)
    assert len(presentation_builder.presentation.slides) == 0

def test_add_slide_adds_slide_to_presentation(presentation_builder, mock_layout):
    slide = presentation_builder.add_slide(mock_layout)
    assert isinstance(slide, Slide)
    assert len(presentation_builder.presentation.slides) == 1

def test_add_slide_returns_slide(presentation_builder, mock_layout):
    slide = presentation_builder.add_slide(mock_layout)
    assert isinstance(slide, Slide)

def test_save_creates_file(presentation_builder, mock_layout, tmpdir):
    # Add a slide to have something in the presentation
    presentation_builder.add_slide(mock_layout)

    # Path to save the file
    file_path = tmpdir.join("test_presentation.pptx")

    # Save and check if the file is created
    presentation_builder.save(str(file_path))
    assert file_path.isfile()



@pytest.fixture
def mock_slide():
    # Create a mock Slide object with necessary attributes
    mock = Mock(spec=Slide)
    mock.shapes = []
    return mock

def test_slide_layout_initialization():
    layout = SlideLayout("Basic", "Arial", "Blue")
    assert layout.layout_name == "Basic"
    assert layout.font == "Arial"
    assert layout.color_scheme == "Blue"

def test_apply_layout_changes_slide(mock_slide):
    # Arrange
    layout = SlideLayout("Title Slide", "Times New Roman", "Red")
    
    # Add mock text frame with paragraph and run with font attribute
    text_shape = Mock()
    text_shape.text = "Example Text"
    paragraph = Mock()
    run = Mock()
    run.font = Mock()
    paragraph.runs = [run]
    text_shape.text_frame.paragraphs = [paragraph]
    mock_slide.shapes.append(text_shape)
    
    # Act
    modified_slide = layout.apply_layout(mock_slide)
    
    # Assert that fonts were changed in mock
    for shape in modified_slide.shapes:
        if hasattr(shape, "text"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    assert run.font.name == "Times New Roman"

def test_apply_layout_returns_slide(mock_slide):
    layout = SlideLayout("Content Slide", "Verdana", "Green")
    result = layout.apply_layout(mock_slide)
    assert isinstance(result, Slide)



@pytest.fixture
def sample_chart():
    # Create a simple Plotly chart
    fig = go.Figure(data=go.Bar(y=[2, 3, 1]))
    return fig

@pytest.fixture
def empty_slide():
    # Create a slide object from a new presentation
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Choosing a blank layout
    return presentation.slides.add_slide(slide_layout)

def test_plotly_chart_embedder_initialization(sample_chart):
    embedder = PlotlyChartEmbedder(sample_chart)
    assert embedder.chart_data == sample_chart

def test_embed_chart_adds_image_to_slide(sample_chart, empty_slide):
    # Arrange
    embedder = PlotlyChartEmbedder(sample_chart)
    initial_shape_count = len(empty_slide.shapes)
    
    # Act
    embedder.embed_chart(empty_slide, (1, 1, 4, 3))  # Example position using inches
    
    # Assert
    assert len(empty_slide.shapes) == initial_shape_count + 1
    new_shape = empty_slide.shapes[-1]
    assert new_shape.shape_type == MSO_SHAPE.PICTURE

def test_embed_chart_correct_position_and_size(sample_chart, empty_slide):
    # Arrange
    embedder = PlotlyChartEmbedder(sample_chart)
    x, y, width, height = 1, 1, 4, 3
    
    # Act
    embedder.embed_chart(empty_slide, (x, y, width, height))
    
    # Assert
    new_shape = empty_slide.shapes[-1]
    assert new_shape.left == Inches(x)
    assert new_shape.top == Inches(y)
    assert new_shape.width == Inches(width)
    assert new_shape.height == Inches(height)



@pytest.fixture
def mock_presentation():
    # Create a mock presentation object with a slide and a text shape
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Using a blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Add a mock text box with text
    text_box = slide.shapes.add_textbox(0, 0, width=Inches(5), height=Inches(1))
    text_frame = text_box.text_frame
    paragraph = text_frame.add_paragraph()
    run = paragraph.add_run()
    run.text = "Sample Text"
    
    return presentation

def test_branding_formatter_initialization():
    guidelines = {"font": "Calibri", "color_scheme": "FF5733"}
    formatter = BrandingFormatter(guidelines)
    assert formatter.branding_guidelines == guidelines

def test_apply_branding_changes_font_and_color(mock_presentation):
    guidelines = {"font": "Calibri", "color_scheme": "FF5733"}
    formatter = BrandingFormatter(guidelines)
    
    # Apply branding
    formatted_presentation = formatter.apply_branding(mock_presentation)
    
    # Assert font and color changes in text runs
    for slide in formatted_presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        assert run.font.name == "Calibri"
                        assert run.font.color.rgb == "FF5733"

def test_apply_branding_applies_to_all_slides(mock_presentation):
    guidelines = {"font": "Times New Roman", "color_scheme": "0000FF"}
    formatter = BrandingFormatter(guidelines)
    
    # Apply branding
    formatted_presentation = formatter.apply_branding(mock_presentation)
    
    # Test functionalities on additional slides
    slide_layout = formatted_presentation.slide_layouts[5]
    additional_slide = formatted_presentation.slides.add_slide(slide_layout)
    
    # Add a new shape to additional slide
    additional_text_box = additional_slide.shapes.add_textbox(0, 0, width=Inches(5), height=Inches(1))
    text_frame = additional_text_box.text_frame
    additional_paragraph = text_frame.add_paragraph()
    additional_run = additional_paragraph.add_run()
    additional_run.text = "Additional Text"
    
    # Assert application of branding
    for slide in formatted_presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        assert run.font.name == "Times New Roman"
                        assert run.font.color.rgb == "0000FF"



def test_load_branding_guidelines_file_not_found():
    with pytest.raises(FileNotFoundError):
        load_branding_guidelines("non_existent_file.json")

def test_load_branding_guidelines_invalid_json():
    invalid_json_content = "{invalid json data}"
    
    # Mock opening a file and raising JSON decode error
    with mock.patch("builtins.open", mock.mock_open(read_data=invalid_json_content)):
        with pytest.raises(ValueError, match="Error parsing the branding guidelines file"):
            load_branding_guidelines("dummy.json")

def test_load_branding_guidelines_valid_json(tmpdir):
    # Create a temporary JSON file with valid content
    data = {
        "font": "Arial",
        "color_scheme": "FF5733"
    }
    temp_file = tmpdir.join("guidelines.json")
    temp_file.write(json.dumps(data))
    
    # Test loading from this temporary JSON file
    result = load_branding_guidelines(str(temp_file))
    
    assert result["font"] == "Arial"
    assert result["color_scheme"] == "FF5733"



# Mock SlideLayout if necessary
class MockSlideLayout(SlideLayout):
    def apply_layout(self, slide):
        pass

@pytest.fixture
def sample_layout_details():
    return {
        'title_position': (1, 0.5, 6, 1),
        'content_position': (1, 2, 6, 3),
        'background_color': 'FFFF00',
        'font_style': {'family': 'Times New Roman', 'size': 18, 'color': 'FF0000'}
    }

@pytest.fixture
def sample_presentation():
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]  # Using blank layout
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Title"
    slide.placeholders[1].text = "Content"
    return slide

def test_create_custom_layout(sample_layout_details, sample_presentation):
    custom_layout_class = create_custom_layout(sample_layout_details)
    assert issubclass(custom_layout_class.__class__, SlideLayout)

    custom_layout = custom_layout_class()
    modified_slide = custom_layout.apply_layout(sample_presentation)
    
    # Verify title position
    title = modified_slide.shapes.title
    assert title.left == Inches(sample_layout_details['title_position'][0])
    assert title.top == Inches(sample_layout_details['title_position'][1])
    
    # Verify content position
    content = modified_slide.placeholders[1]
    assert content.left == Inches(sample_layout_details['content_position'][0])
    assert content.top == Inches(sample_layout_details['content_position'][1])

    # Verify background color
    background_color_rgb = RGBColor.from_string(sample_layout_details['background_color'])
    assert modified_slide.background.fill.fore_color.rgb == background_color_rgb

    # Verify font style
    font_family = sample_layout_details['font_style']['family']
    font_size = sample_layout_details['font_style']['size']
    font_color_rgb = RGBColor.from_string(sample_layout_details['font_style']['color'])
    
    for paragraph in title.text_frame.paragraphs:
        assert paragraph.font.size == font_size
        assert paragraph.font.name == font_family
        assert paragraph.font.color.rgb == font_color_rgb

def test_create_custom_layout_default_values():
    custom_layout_class = create_custom_layout({})
    assert issubclass(custom_layout_class.__class__, SlideLayout)

    # Create a mock slide to test defaults
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Default Title"
    
    custom_layout = custom_layout_class()
    modified_slide = custom_layout.apply_layout(slide)

    # Verify default title position
    default_title_position = (1, 1, 5, 1)
    assert modified_slide.shapes.title.left == Inches(default_title_position[0])
    assert modified_slide.shapes.title.top == Inches(default_title_position[1])



def test_generate_bar_chart_with_dataframe():
    df = pd.DataFrame({'A': [1, 2, 3], 'B': [4, 5, 6]})
    fig = generate_interactive_chart(df, 'bar')
    assert isinstance(fig, Figure)
    assert len(fig.data) == 1
    assert fig.data[0].type == 'bar'

def test_generate_bar_chart_with_dict():
    data = {'A': 1, 'B': 2, 'C': 3}
    fig = generate_interactive_chart(data, 'bar')
    assert isinstance(fig, Figure)
    assert len(fig.data) == 1
    assert fig.data[0].type == 'bar'

def test_generate_line_chart_with_dataframe():
    df = pd.DataFrame({'X': [1, 2, 3], 'Y': [4, 5, 6]})
    fig = generate_interactive_chart(df, 'line')
    assert isinstance(fig, Figure)
    assert len(fig.data) == 1
    assert fig.data[0].type == 'scatter'
    assert fig.data[0].mode == 'lines'

def test_generate_pie_chart_with_dict():
    data = {'Category 1': 10, 'Category 2': 20, 'Category 3': 30}
    fig = generate_interactive_chart(data, 'pie')
    assert isinstance(fig, Figure)
    assert len(fig.data) == 1
    assert fig.data[0].type == 'pie'

def test_invalid_chart_type():
    df = pd.DataFrame({'X': [1, 2, 3], 'Y': [4, 5, 6]})
    with pytest.raises(ValueError, match="Unsupported chart type"):
        generate_interactive_chart(df, 'invalid_type')

def test_invalid_data_format():
    with pytest.raises(ValueError, match="Data should be a pandas DataFrame, dictionary, or list."):
        generate_interactive_chart("invalid", 'bar')
