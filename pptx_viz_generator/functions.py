
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_LABEL_POSITION, XL_TICK_LABEL_POSITION, XL_LEGEND_POSITION, XL_CHART_TYPE
import plotly.graph_objects as go
import os

def create_presentation():
    """Create a new PowerPoint presentation object."""
    return Presentation()

def add_title_slide(prs, title):
    """Add a title slide to the presentation."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

def add_slide(presentation, layout):
    """Add a slide with a specific layout to the presentation."""
    slide_layout = presentation.slide_layouts[layout]
    return presentation.slides.add_slide(slide_layout)

def add_text_to_slide(presentation, slide_index, text):
    """Add text to a slide."""
    slide = presentation.slides[slide_index]
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
    textbox.text_frame.text = text

def add_image_to_slide(slide, image_path, left, top, width, height):
    """Add an image to a slide."""
    slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))

def add_plotly_chart_to_slide(prs, slide_index, chart):
    """Add a Plotly chart to a slide."""
    slide = prs.slides[slide_index]
    left, top, width, height = Inches(1), Inches(1), Inches(6), Inches(4.5)
    
    chart_image_path = "chart.png"
    chart.write_image(chart_image_path)
    
    slide.shapes.add_picture(chart_image_path, left, top, width, height)
    
    os.remove(chart_image_path)

def format_text_font(text, font_name, font_size, bold=False, italic=False, underline=False, align=None):
    """Format the font of text on a slide."""
    text.font.name = font_name
    text.font.size = Pt(font_size)
    text.font.bold = bold
    text.font.italic = italic
    text.font.underline = underline
    
    if align == 'left':
        text.alignment = PP_ALIGN.LEFT
    elif align == 'center':
        text.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        text.alignment = PP_ALIGN.RIGHT

def format_background_color(slide, color):
    """Formats the background color of a slide."""
    rgb_color = RGBColor(color[0], color[1], color[2])
    
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

def format_image(image_path, slide_number, left_px=0.0 , top_px=0.0 , width_px=None , height_px=None ):
    
     prs= Presentation()
     # Get the specified slide 
     sld=prs.slides.add_slide(prs.slide_layouts[slide_number])
     
     # Add image to the side 
     sld.shapes.add_picture(image_path,left=Inches(left_px) ,top=Inches(top_px) ,width=Inches(width_px) ,height=Inches(height_px))
     
     prs.save('formatted_presentation.pptx')


def format_chart_lines(chart ,line_style="solid" , line_color="000000"):
   
   for trace in chart.data:
       trace.line=dict(dash=line_style)
       trace.line.color=line_color



def format_data_labels(chart):
     for series in chart.series:
         series.has_data_labels=True 
         data_labels=series.data_labels

         data_labels.position=XL_LABEL_POSITION.OUTSIDE_END

         font=data_labels.font 
         font.color.rgb=RGBColor(255 , 0 , 0) 
         font.size=Pt(12)


def format_axis_labels(chart,font_name="Arial", font_size=None ,font_color=None):

   x_axis,y_axis=chart.category_axis , chart.value_axis
   
   if font_name:
       x_axis.tick_labels.font.name=y_axis.tick_labels.font.name=font_name

   if font_size:
       x_axis.tick_labels.font.size=y_axis.tick_labels.font.size=font_size 

   if font_color:
      x_axis.tick_labels.font.color.rgb=y_axis.tick_labels.font.color.rgb=font_color


   x_axis.tick_label_position=y_axis.tick_label_position=XL_TICK_LABEL_POSITION.LOW


def format_chart_legend(chart , position=None ,font_color=None):

     if position or font_color:
          chart.has_legend=True 

          if position:
              chart.legend.position=position

          if font_color:
              for legend in chart.legend.legend_entries: legend.font.color.rgb=font_color


def export_presentation(presentation , output_file): presentation.save(output_file)


def save_ppt_as_pdf(presentation_path,pdf_path):

     prsentation.save(presentation_path).save(pdf_path,"PDF")




#Set default properties for PowerPoint slides  

def set_default_font(presentation,font_name):

  for sld in presentation.slides:

      for shpe in sld.shapes:

            if shpe.has_text_frame:

                for prgph in shpe.text_frame.paragraphs: 
                    
                    for rn in prgph.runs: rn.font.name=font_name



def set_default_font_size(presentation,font_size):

  for sld in presentation.slides:

      for shpe in sld.shapes:

           if shpe.has_text_frame:
            
              for prgph in shpe.text_frame.paragraphs:

                  for rn in prgph.runs:  rn.font.size=Pt(font_size)


     

def set_default_font_color(presentation,font_color):

  for sld in presentation.slides:

      for shpe in sld.shapes:

           if shpe.has_text_frame:
            
              for prgph in shpe.text_frame.paragraphs:

                  for rn in prgph.runs:  rn.font.color.rgb=str(font_color)



#Set default background color of slides 

def set_background_color(presentation,color):

   rgb_clr=str(color)
   
   for sld in presentation.slides:

         fill=sld.background.fill 
        
         fill.solid() 
        
         fill.fore_color.rgb=str(rgb_clr)

        

#Set default line style and color of charts  

def set_default_line_style(presentaion,line_style,line_hex_code):
       
       line_clr=str(line_hex_code)

       for sld in presentaion.slides:

             for shape in sld.shapes : 

                   if shape.has_chart : 

                           chart_type_=shape.chart.chart_type 
                           
                           #Check whether it is line type or not  
                           is_line_or_not=lambda chrt_type : chrt_type==XL_CHART_TYPE.LINE 

                           #Assign default line style and color to it   
                           if is_line_or_not(chart_type_) :

                                shape.chart.plots[0].format.line.color.rgb=str(line_clr)
                                
                                shape.chart.plots[0].format.line.dash=line_style




#Set default marker style and marker color  

def set_default_marker_style_and_color(presentaion,mk_style,mk_hex_code):

      mk_clr=str(mk_hex_code)

      def assign_marker_sty_and_clr(series_,style_,color_):
             
             series_.marker.fill.solid()
             
             series_.marker.style=str(style_)
             
             series_.marker.fill.fore_color.rgb=color_

      
      for sld_ in presentaion.slides:

           for shape_ in sld_.shapes :    

               if shape_.has_chart :     

                    assign_marker_sty_and_clr(shape_.chart.series,mk_style,str(mk_clr))



#Applying branding guidelines and generate multiple slides containing data visualizations from input data  
#
      
      
      
      
      
      
      

 def apply_custom_branding_and_generate_slides(prsn,data_lst,bckgrnd_rgb_colrs_dict,fnt_rgb_colrs_dict,fnt_nms_dict,fnt_sz_lst,max_fnt_sz,max_chrt_sz,mrk_styl_lst,mrk_hex_cd_lst,line_styl_lst,line_hex_cd_lst):

     
        def set_branding_guidelines(font_colors_:dict,bckgrnd_rgb_colors_:dict,font_names_:dict,font_sizes_:list,max_font_sizes:int,max_chart_sizes:int,
                                                      marker_styles:list=[],marker_colors:list=[],line_styles:list=[],line_colors:list=[]): return dict(font_colors_=font_colors_,
                                                                                                                                                                  bckgrnd_rgb_colors_=bckgrnd_rgb_colors_,font_names_=font_names_,font_sizes_=font_sizes,
                                                                                                                                                                  max_font_sizes=max_font_sizes,max_chart_sizes=max_chart_sizes,
                                                                                                                                                                  marker_styles_=marker_styles_,marker_colors_=marker_colors_,
                                                                                                                                                                  line_styles_=line_styles_,line_colors_=line_colors_)
        
        
        def apply_branding_for_each_slide_in_presentaion(slide_obj,bckgrnd_rgb_colrs_dict,fnt_rgb_colrs_dict,fnt_nms_dict,fnt_sz_lst):
            
            #Apply branding guidelines on each created new slides by user
            
                apply_branding(slide_obj,set_branding_guidelines(fnt_rgb_colrs_dict,bckgrnd_rgb_colrs_dict,fnt_nms_dict,fnt_sz_lst,max_fnt_sz,max_chrt_sz,mrk_styl_lst,mrk_hex_cd_lst,line_styl_lst,line_hex_cd_lst))

                
                apply_custom_branding_and_generate_slides(create_presentation(),data_lst,bckgrnd_rgb_colrs_dict,fnt_rgb_colrs_dict,fnt_nms_dict,fnt_sz_lst,
                                                                                                              max_fnt_sz,max_chrt_sz,mrk_styl_lst,mrk_hex_cd_lst,line_styl_lst,line_hex_cd_lst) 


        def generate_slides(data,presentation):  

                 apply_branding_for_each_slide_in_presentaion(set_background(presentation,bckgrnd_rgb_colrs_dict),
                                                                                                          fnt_rgb_colrs_dict,fnt_nms_dict,[None])

                 #Create new slides containing data visualizations from input data 
                 
                 ##Loop through each data sets and create new slides containing visualizations based on Plotly library
                
                
                 def convert_plotly_figures_to_imgs(data): return [go.Figure(data).write_image("temp_{i}.png")  and f"temp_{i}.png".remove()]

                 
                 ##Loop through each data sets and create new slides containing visualizations based on Plotly library
                
                 return [add_plotly_chart_to_slide(add_title_slide(generate_slides(data,presentation),'Slide Title Here')
                                                                        ,"Slide Subtitle Here",convert_plotly_figures_to_imgs(data(i)) )  and os.remove(f"temp_{i}.png") ] 



